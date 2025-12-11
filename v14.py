import os
import io
import uuid
import tempfile
import base64
import time
from datetime import datetime, timedelta
import difflib
import math

import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st

# Optional: for PPTX export (if not installed, function falls back to PDF)
try:
    from pptx import Presentation
    from pptx.util import Inches

    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

# ==============================================================================
# AURORA DASHBOARD v8.0 ULTIMATE (ENTERPRISE EDITION)
# ==============================================================================
# A single-file, production-grade Streamlit application featuring:
# - Advanced Financial Analytics (Monte Carlo, FFT, Volatility Cones)
# - Real-time Trading Simulation & Order Book
# - Customer Cohort & RFM Analysis
# - PyQt5-based High-Fidelity PDF Reporting Engine (Subprocess Architecture)
# - Adaptive Glassmorphism UI (The "Aurora" Design System)
# - Robust Session Management & Data Caching
#
# AUTHOR: Gemini (Refined for Enterprise Scale)
# VERSION: 8.0.1
# ==============================================================================

from matplotlib.pyplot import title
import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime, timedelta, time as dt_time
import io
import base64
import json
import math
import random
import textwrap
import sys
import os
import tempfile
import subprocess
import uuid
import calendar
import threading
import smtplib
import time
from pathlib import Path
from email.message import EmailMessage
from email.utils import formatdate
import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime, timedelta
import time
import random
import math
import io
import base64

# Scientific Computing Imports
try:
    import numpy.fft as fft
    from scipy import stats
except ImportError:
    st.error("Scientific libraries (scipy) missing. Please install: pip install scipy")

# ==============================================================================
# 1. GLOBAL CONFIGURATION & CONSTANTS
# ==============================================================================

# New Imports for Reporting
from fpdf import FPDF

# Page Setup (Keep existing)
st.set_page_config(
    page_title="Aurora V7 | Ultimate Dashboard",
    layout="wide",
    page_icon="💠",
    initial_sidebar_state="expanded",
)

# Constants for simulations
TRADING_DAYS = 252
RISK_FREE_RATE = 0.02
SIMULATION_RUNS = 1000

# ==============================================================================
# 2. SESSION STATE MANAGEMENT
# ==============================================================================


def init_session_state():
    """Initializes all session state variables with default values."""
    defaults = {
        "theme": "Revolut Space Blue",
        "layout_density": "Comfortable",
        "user": {
            "id": f"usr_{uuid.uuid4().hex[:8]}",
            "name": "Enterprise Admin",
            "email": "admin@aurora-finance.io",
            "avatar": "https://ui-avatars.com/api/?name=Admin&background=00d4ff&color=fff",
            "role": "Chief Investment Officer",
            "bio": "Overseeing global asset allocation and risk management.",
        },
        "notifications": [
            {
                "id": 101,
                "title": "Market Volatility Alert",
                "msg": "VIX index spiked by 12% in the last hour.",
                "level": "warning",
                "time": "1h ago",
            },
            {
                "id": 102,
                "title": "Portfolio Rebalancing",
                "msg": "Auto-rebalancing completed successfully.",
                "level": "success",
                "time": "4h ago",
            },
            {
                "id": 103,
                "title": "System Update",
                "msg": "Aurora v8.0 patch applied.",
                "level": "info",
                "time": "1d ago",
            },
        ],
        "activity_log": [
            {
                "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                "user": "System",
                "action": "Dashboard Initialized",
                "details": "v8.0 Boot Sequence",
            }
        ],
        "smtp_config": {
            "host": "smtp.example.com",
            "port": 587,
            "user": "",
            "pass": "",
            "from": "reports@aurora.io",
        },
    }

    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

    if "activity_log" not in st.session_state:
        st.session_state.activity_log = []
    if "comments" not in st.session_state:
        st.session_state.comments = {}  # { context_id: [ {user, time, text} ] }
    if "saved_views" not in st.session_state:
        st.session_state.saved_views = {}  # { view_name: { config_dict } }

    if "activity_log" not in st.session_state:
        st.session_state.activity_log = []
    if "comments" not in st.session_state:
        st.session_state.comments = {}  # { context_id: [ {user, time, text} ] }


init_session_state()

# ==============================================================================
# 3. DATA GENERATION ENGINE (ADVANCED)
# ==============================================================================


class DataEngine:
    """
    Centralized data generation and management class.
    Uses caching to prevent re-computation on every rerun.
    """

    @staticmethod
    @st.cache_data(ttl=3600, show_spinner=False)
    def generate_market_data(days=500):
        """Generates complex market data with geometric brownian motion properties."""
        np.random.seed(42)
        dates = pd.date_range(end=datetime.now(), periods=days)

        # Parameters for synthetic asset
        dt = 1 / TRADING_DAYS
        mu = 0.08  # Drift
        sigma = 0.2  # Volatility

        # Geometric Brownian Motion
        returns = np.random.normal(
            (mu - 0.5 * sigma**2) * dt, sigma * np.sqrt(dt), days
        )
        price_path = 100 * np.exp(np.cumsum(returns))

        # Volume with lognormal distribution + regime spikes
        volume = np.random.lognormal(10, 0.8, days)

        # Moving Averages
        df = pd.DataFrame({"date": dates, "price": price_path, "volume": volume})
        df["returns_pct"] = df["price"].pct_change().fillna(0)
        df["ma_20"] = df["price"].rolling(20).mean()
        df["ma_50"] = df["price"].rolling(50).mean()
        df["ma_200"] = df["price"].rolling(200).mean()
        df["volatility_20d"] = df["returns_pct"].rolling(20).std() * np.sqrt(252)

        # RSI Calculation
        delta = df["price"].diff()
        gain = (delta.where(delta > 0, 0)).rolling(14).mean()
        loss = (-delta.where(delta < 0, 0)).rolling(14).mean()
        rs = gain / loss
        df["rsi"] = 100 - (100 / (1 + rs))

        return df

    @staticmethod
    @st.cache_data(ttl=3600, show_spinner=False)
    def generate_sales_data(records=2000):
        """Generates detailed transactional sales data for cohort analysis."""
        np.random.seed(101)
        dates = pd.date_range(start="2023-01-01", end="2024-12-31", freq="H")

        # Sampling random timestamps
        sample_dates = np.random.choice(dates, records)

        products = [
            "Aurora Premium",
            "Aurora Lite",
            "Data Add-on",
            "Consulting Hour",
            "API Access",
        ]
        regions = ["North America", "EMEA", "APAC", "LATAM"]
        channels = ["Direct Sales", "Partner", "Web Organic", "Referral"]

        data = {
            "transaction_id": [
                f"TRX-{uuid.uuid4().hex[:8].upper()}" for _ in range(records)
            ],
            "date": sample_dates,
            "product": np.random.choice(products, records, p=[0.2, 0.4, 0.2, 0.1, 0.1]),
            "region": np.random.choice(regions, records, p=[0.4, 0.3, 0.2, 0.1]),
            "channel": np.random.choice(channels, records),
            "amount": np.random.lognormal(
                5, 1, records
            ),  # Log-normal for realistic pricing
            "cost": np.zeros(records),
            "customer_id": np.random.choice(
                [f"CUST-{i:04d}" for i in range(1, 501)], records
            ),
        }

        df = pd.DataFrame(data)
        df["date"] = pd.to_datetime(df["date"])
        df["month"] = df["date"].dt.to_period("M")
        df["cost"] = df["amount"] * np.random.uniform(0.3, 0.6, records)
        df["profit"] = df["amount"] - df["cost"]
        df["margin"] = df["profit"] / df["amount"]

        # Simulate Customer Acquisition Date for Cohorts
        cust_signup = df.groupby("customer_id")["date"].min().to_dict()
        df["signup_date"] = df["customer_id"].map(cust_signup)
        df["cohort_month"] = df["signup_date"].dt.to_period("M")

        return df.sort_values("date")

    @staticmethod
    @st.cache_data(ttl=3600, show_spinner=False)
    def generate_portfolio_data():
        """Generates a multi-asset portfolio with sector allocation."""
        assets = [
            {
                "ticker": "AAPL",
                "name": "Apple Inc.",
                "sector": "Tech",
                "type": "Equity",
            },
            {
                "ticker": "NVDA",
                "name": "Nvidia Corp.",
                "sector": "Tech",
                "type": "Equity",
            },
            {
                "ticker": "JPM",
                "name": "JPMorgan",
                "sector": "Finance",
                "type": "Equity",
            },
            {
                "ticker": "XOM",
                "name": "Exxon Mobil",
                "sector": "Energy",
                "type": "Equity",
            },
            {
                "ticker": "BND",
                "name": "Total Bond Mkt",
                "sector": "Fixed Income",
                "type": "ETF",
            },
            {
                "ticker": "GLD",
                "name": "SPDR Gold",
                "sector": "Commodity",
                "type": "ETF",
            },
            {
                "ticker": "BTC-USD",
                "name": "Bitcoin",
                "sector": "Crypto",
                "type": "Crypto",
            },
        ]

        data = []
        for asset in assets:
            qty = np.random.randint(50, 5000)
            price = np.random.uniform(100, 1000)
            data.append(
                {
                    **asset,
                    "quantity": qty,
                    "avg_price": price
                    * np.random.uniform(0.8, 1.1),  # Current price vs avg cost
                    "current_price": price,
                    "daily_change_pct": np.random.normal(0, 0.02),
                }
            )

        df = pd.DataFrame(data)
        df["market_value"] = df["quantity"] * df["current_price"]
        df["cost_basis"] = df["quantity"] * df["avg_price"]
        df["unrealized_pl"] = df["market_value"] - df["cost_basis"]
        df["unrealized_pl_pct"] = (df["unrealized_pl"] / df["cost_basis"]) * 100
        df["weight_pct"] = df["market_value"] / df["market_value"].sum() * 100

        return df

    @staticmethod
    @st.cache_data(ttl=3600, show_spinner=False)
    def generate_company_registry(records=500):
        """Generates a registry of companies for the Search Engine."""
        np.random.seed(2024)
        sectors = [
            "Technology",
            "Healthcare",
            "Finance",
            "Energy",
            "Consumer",
            "Industrial",
        ]
        status_opts = ["Active", "Inactive", "Pending", "Merged"]

        data = {
            "company_id": [f"COMP-{1000 + i}" for i in range(records)],
            "company_name": [
                f"Aurora Corp {i}" for i in range(records)
            ],  # Placeholder names
            "sector": np.random.choice(sectors, records),
            "founded_year": np.random.randint(1950, 2023, records),
            "status": np.random.choice(status_opts, records, p=[0.8, 0.1, 0.05, 0.05]),
            "parent_id": [
                f"COMP-{np.random.randint(1000, 1050)}"
                if i > 50 and random.random() > 0.7
                else None
                for i in range(records)
            ],
            "subsidiaries_count": np.random.randint(0, 15, records),
            "revenue_mm": np.random.uniform(10, 5000, records),
            "employees": np.random.randint(10, 50000, records),
            "risk_score": np.random.randint(1, 100, records),
        }

        df = pd.DataFrame(data)
        # Add some specific names for search demo
        df.loc[0, "company_name"] = "Acme Corp International"
        df.loc[1, "company_name"] = "Wayne Enterprises"
        df.loc[2, "company_name"] = "Stark Industries"
        df.loc[3, "company_name"] = "Cyberdyne Systems"

        return df


# Initialize Data in Session State
if "market_df" not in st.session_state:
    st.session_state.market_df = DataEngine.generate_market_data()
if "sales_df" not in st.session_state:
    st.session_state.sales_df = DataEngine.generate_sales_data()
if "portfolio_df" not in st.session_state:
    st.session_state.portfolio_df = DataEngine.generate_portfolio_data()
if "company_df" not in st.session_state:
    st.session_state.company_df = DataEngine.generate_company_registry()

# ==============================================================================
# 4. THEME ENGINE & UI SYSTEM
# ==============================================================================


class ThemeEngine:
    """Handles CSS injection and theme management."""

    THEMES = {
        "Revolut Space Blue": {
            "bg": "#071028",
            "card": "rgba(255,255,255,0.03)",
            "glass": "rgba(255,255,255,0.04)",
            "primary": "#00d4ff",
            "accent": "#7b2ff7",
            "text": "#ffffff",
            "muted": "#9fb0c8",
            "gradient": "linear-gradient(135deg, #00d4ff 0%, #7b2ff7 100%)",
            "success": "#00e396",
            "danger": "#ff0055",
            "warning": "#feb019",
        },
        "Cyberpunk Neon": {
            "bg": "#050505",
            "card": "rgba(20,20,20,0.8)",
            "glass": "rgba(40,40,40,0.5)",
            "primary": "#fcee0a",
            "accent": "#ff003c",
            "text": "#e0e0e0",
            "muted": "#666",
            "gradient": "linear-gradient(135deg, #fcee0a 0%, #ff003c 100%)",
            "success": "#0aff0a",
            "danger": "#ff003c",
            "warning": "#ffaa00",
        },
        "Corporate Slate": {
            "bg": "#f0f2f6",
            "card": "#ffffff",
            "glass": "rgba(255,255,255,0.9)",
            "primary": "#2b3a55",
            "accent": "#ce7777",
            "text": "#1a1a1a",
            "muted": "#888888",
            "gradient": "linear-gradient(135deg, #2b3a55 0%, #4a5568 100%)",
            "success": "#38a169",
            "danger": "#e53e3e",
            "warning": "#d69e2e",
        },
    }

    @classmethod
    def apply_theme(cls):
        theme_name = st.session_state.theme
        t = cls.THEMES.get(theme_name, cls.THEMES["Revolut Space Blue"])

        # Conditional font color for light/dark themes
        is_light = theme_name == "Corporate Slate"
        text_color = "#1a1a1a" if is_light else "#ffffff"

        css = f"""
        <style>
            @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;700&display=swap');
            
            :root {{
                --bg-color: {t["bg"]};
                --card-color: {t["card"]};
                --glass-color: {t["glass"]};
                --primary-color: {t["primary"]};
                --accent-color: {t["accent"]};
                --text-color: {text_color};
                --muted-color: {t["muted"]};
                --gradient: {t["gradient"]};
                --success: {t["success"]};
                --danger: {t["danger"]};
            }}

            /* Base App Styling */
            .stApp {{
                background-color: var(--bg-color);
                color: var(--text-color);
                font-family: 'Inter', sans-serif;
            }}
            
            h1, h2, h3, h4, h5, h6 {{
                font-family: 'Inter', sans-serif;
                font-weight: 700;
                color: var(--text-color) !important;
            }}
            
            /* Navigation Bar */
            .nav-container {{
                position: sticky; top: 0; z-index: 1000;
                background: var(--glass-color);
                backdrop-filter: blur(12px); -webkit-backdrop-filter: blur(12px);
                border-bottom: 1px solid rgba(255,255,255,0.05);
                padding: 1rem 1.5rem;
                display: flex; align-items: center; justify-content: space-between;
                margin-bottom: 2rem;
                border-radius: 0 0 16px 16px;
                box-shadow: 0 4px 20px rgba(0,0,0,0.1);
            }}
            
            .nav-brand {{
                font-size: 1.5rem; font-weight: 800;
                background: var(--gradient);
                -webkit-background-clip: text; -webkit-text-fill-color: transparent;
                letter-spacing: -0.5px;
            }}

            /* Glassmorphism Cards */
            .glass-card {{
                background: var(--glass-color);
                border: 1px solid rgba(255,255,255,0.05);
                border-radius: 16px;
                padding: 1.5rem;
                box-shadow: 0 8px 32px rgba(0,0,0,0.15);
                transition: transform 0.2s ease, border-color 0.2s ease;
                backdrop-filter: blur(8px);
            }}
            .glass-card:hover {{
                transform: translateY(-2px);
                border-color: rgba(255,255,255,0.15);
            }}

            /* KPI Cards */
            .kpi-metric {{
                text-align: center;
                padding: 1rem;
                border-radius: 12px;
                background: linear-gradient(145deg, rgba(255,255,255,0.03), rgba(255,255,255,0.01));
                border: 1px solid rgba(255,255,255,0.05);
            }}
            .kpi-value {{
                font-size: 1.8rem; font-weight: 700;
                color: var(--primary-color);
                font-family: 'JetBrains Mono', monospace;
            }}
            .kpi-label {{
                font-size: 0.85rem; color: var(--muted-color);
                text-transform: uppercase; letter-spacing: 1px;
                margin-bottom: 0.5rem;
            }}
            .kpi-delta {{
                font-size: 0.9rem; font-weight: 600;
            }}
            .kpi-delta.positive {{ color: var(--success); }}
            .kpi-delta.negative {{ color: var(--danger); }}

            /* Buttons & Inputs */
            .stButton > button {{
                background: var(--gradient);
                border: none;
                color: #fff;
                font-weight: 600;
                padding: 0.6rem 1.2rem;
                border-radius: 8px;
                transition: opacity 0.2s;
            }}
            .stButton > button:hover {{ opacity: 0.9; }}
            
            /* Sidebar */
            [data-testid="stSidebar"] {{
                background-color: rgba(0,0,0,0.2);
                border-right: 1px solid rgba(255,255,255,0.05);
            }}
            
            /* Scrollbars */
            ::-webkit-scrollbar {{ width: 8px; height: 8px; }}
            ::-webkit-scrollbar-track {{ background: transparent; }}
            ::-webkit-scrollbar-thumb {{ background: rgba(255,255,255,0.1); border-radius: 4px; }}
            ::-webkit-scrollbar-thumb:hover {{ background: rgba(255,255,255,0.2); }}
            
            /* Utility Classes */
            .text-muted {{ color: var(--muted-color); }}
            .text-primary {{ color: var(--primary-color); }}
            .text-success {{ color: var(--success); }}
            .text-danger {{ color: var(--danger); }}
        </style>
        """
        st.markdown(css, unsafe_allow_html=True)

    @staticmethod
    def render_header(title, subtitle=None):
        sub_html = (
            f"<div style='color:var(--muted-color); font-size:0.9rem; margin-top:-5px'>{subtitle}</div>"
            if subtitle
            else ""
        )
        st.markdown(
            f"""
        <div class="nav-container">
            <div>
                <div class="nav-brand">{title}</div>
                {sub_html}
            </div>
            <div style="display:flex; gap:10px; align-items:center;">
                <div style="text-align:right;">
                    <div style="font-weight:600; font-size:0.9rem;">{st.session_state.user["name"]}</div>
                    <div style="font-size:0.75rem; color:var(--muted-color);">{st.session_state.user["role"]}</div>
                </div>
                <img src="{st.session_state.user["avatar"]}" style="width:40px; height:40px; border-radius:50%; border:2px solid var(--primary-color);">
            </div>
        </div>
        """,
            unsafe_allow_html=True,
        )

    @staticmethod
    def render_kpi_card(label, value, delta=None, delta_desc="vs last period"):
        delta_html = ""
        if delta:
            color_cls = (
                "positive"
                if "+" in delta or float(delta.strip("%+")) >= 0
                else "negative"
            )
            delta_html = f"<div class='kpi-delta {color_cls}'>{delta} <span style='font-size:0.7em; color:var(--muted-color); font-weight:400'>{delta_desc}</span></div>"

        return f"""
        <div class="kpi-metric">
            <div class="kpi-label">{label}</div>
            <div class="kpi-value">{value}</div>
            {delta_html}
        </div>
        """

    @staticmethod
    def render_comment_section(context_id):
        """Renders a commenting widget for a specific context (page or entity)."""
        st.markdown("---")
        st.markdown(f"#### 💬 Discussion & Notes")

        # Display existing comments
        if context_id in st.session_state.comments:
            for c in st.session_state.comments[context_id]:
                st.markdown(
                    f"""
                <div style="background:rgba(255,255,255,0.03); padding:10px; border-radius:8px; margin-bottom:10px; border-left: 3px solid var(--primary-color)">
                    <div style="display:flex; justify-content:space-between; font-size:0.75rem; color:var(--muted-color)">
                        <span>{c["user"]}</span>
                        <span>{c["time"]}</span>
                    </div>
                    <div style="margin-top:5px; font-size:0.9rem;">{c["text"]}</div>
                </div>
                """,
                    unsafe_allow_html=True,
                )
        else:
            st.info("No comments yet. Be the first to start the discussion.")

        # Add new comment
        with st.form(key=f"comment_form_{context_id}"):
            new_comment = st.text_area(
                "Add a comment...", height=68, key=f"txt_{context_id}"
            )
            if st.form_submit_button(
                "Post Comment", type="primary", use_container_width=True
            ):
                if new_comment:
                    if context_id not in st.session_state.comments:
                        st.session_state.comments[context_id] = []

                    st.session_state.comments[context_id].append(
                        {
                            "user": st.session_state.user["name"],
                            "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                            "text": new_comment,
                        }
                    )
                    st.success("Comment posted!")
                    st.rerun()


# Apply the theme immediately
ThemeEngine.apply_theme()

# ==============================================================================
# 5. ADVANCED ANALYTICS LIBRARY
# ==============================================================================


class AnalyticsLib:
    """Library of heavy-duty analytical functions."""

    @staticmethod
    def monte_carlo_simulation(df, days=30, simulations=1000):
        """
        Runs Monte Carlo simulation on the provided price series.
        Returns: Simulation DataFrame paths.
        """
        last_price = df["price"].iloc[-1]
        returns = df["returns_pct"]
        volatility = returns.std()

        simulation_df = pd.DataFrame()

        for x in range(simulations):
            count = 0
            price_series = []
            price = last_price * (1 + np.random.normal(0, volatility))
            price_series.append(price)

            for y in range(days):
                if count == 251:
                    break
                price = price_series[count] * (1 + np.random.normal(0, volatility))
                price_series.append(price)
                count += 1

            simulation_df[x] = price_series

        return simulation_df

    @staticmethod
    def calculate_volatility_cones(df, windows=[10, 30, 60, 90, 180]):
        """
        Calculates realized volatility cones for different lookback windows.
        Useful for options pricing and risk management.
        """
        cones = {}
        for w in windows:
            # Annualized volatility
            vol = df["returns_pct"].rolling(w).std() * np.sqrt(252)
            # Get min, max, median, 25th, 75th percentiles of the historical distribution
            cones[w] = {
                "min": vol.min(),
                "max": vol.max(),
                "median": vol.median(),
                "p25": vol.quantile(0.25),
                "p75": vol.quantile(0.75),
                "current": vol.iloc[-1],
            }
        return pd.DataFrame(cones).T

    @staticmethod
    def perform_fft_analysis(df, column="price"):
        """
        Fast Fourier Transform to identify dominant cyclical periods in data.
        """
        data = df[column].values
        # Detrend
        data_detrended = data - np.mean(data)
        n = len(data)
        fhat = fft.fft(data_detrended, n)
        PSD = fhat * np.conj(fhat) / n  # Power Spectral Density
        freq = (1 / (1 * n)) * np.arange(n)
        L = np.arange(1, np.floor(n / 2), dtype="int")

        period = 1 / freq[L]
        power = PSD[L].real

        return pd.DataFrame({"Period (Days)": period, "Power": power}).sort_values(
            "Power", ascending=False
        )


# ==============================================================================
# 6. PDF GENERATION ENGINE (SUBPROCESS ARCHITECTURE)
# ==============================================================================


class PDFEngine:
    """
    Robust PDF generation using a dedicated subprocess script.
    This avoids threading conflicts between Streamlit and PyQt5/QtWebEngine.
    """

    @staticmethod
    def _generate_script(html_path, output_path):
        """Creates the isolated python script content for PDF generation."""
        # Escaping paths for Windows compatibility
        html_path = html_path.replace("\\", "\\\\")
        output_path = output_path.replace("\\", "\\\\")

        return f"""
import sys
import os
from PyQt5.QtWidgets import QApplication
from PyQt5.QtWebEngineWidgets import QWebEnginePage
from PyQt5.QtCore import QUrl, QMarginsF

# --- PDF Handler ---
def pdf_finished(success):
    if success:
        print("STATUS: SUCCESS")
    else:
        print("STATUS: FAILED")
    QApplication.instance().quit()

def main():
    app = QApplication(sys.argv)
    page = QWebEnginePage()
    
    # Load HTML
    url = QUrl.fromLocalFile(r"{html_path}")
    page.load(url)
    
    # Callback
    def on_load(ok):
        if not ok:
            print("STATUS: LOAD_ERROR")
            app.quit()
            return
        # Print options
        page.printToPdf(r"{output_path}")
        
    page.loadFinished.connect(on_load)
    page.pdfPrintingFinished.connect(pdf_finished)
    
    sys.exit(app.exec_())

if __name__ == "__main__":
    main()
"""

    @classmethod
    def generate(cls, html_content, filename="report.pdf"):
        """
        Public method to generate a PDF from HTML content.
        Returns: (success: bool, data_or_error: bytes|str)
        """
        # 1. Write HTML to temporary file
        fd, html_path = tempfile.mkstemp(suffix=".html")
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            f.write(html_content)

        # 2. Define output path
        output_path = os.path.join(
            tempfile.gettempdir(), f"aurora_{uuid.uuid4().hex}.pdf"
        )

        # 3. Create the generator script
        script_content = cls._generate_script(html_path, output_path)
        script_fd, script_path = tempfile.mkstemp(suffix=".py")
        with os.fdopen(script_fd, "w", encoding="utf-8") as f:
            f.write(script_content)

        try:
            # 4. Execute subprocess
            # We use the same python interpreter that is running Streamlit
            cmd = [sys.executable, script_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)

            # 5. Check result
            if "STATUS: SUCCESS" in result.stdout and os.path.exists(output_path):
                with open(output_path, "rb") as f:
                    pdf_data = f.read()
                return True, pdf_data
            else:
                return (
                    False,
                    f"Subprocess Error:\nSTDOUT: {result.stdout}\nSTDERR: {result.stderr}",
                )

        except Exception as e:
            return False, str(e)

        finally:
            # Cleanup temp files
            for p in [html_path, script_path, output_path]:
                if os.path.exists(p):
                    try:
                        os.remove(p)
                    except:
                        pass


# ==============================================================================
# 7. PAGE IMPLEMENTATIONS
# ==============================================================================


# --- DASHBOARD PAGE ---
def page_dashboard():
    ThemeEngine.render_header("Aurora Analytics", "Executive Overview")

    market_df = st.session_state.market_df
    portfolio_df = st.session_state.portfolio_df

    # 1. High-Level KPIs
    c1, c2, c3, c4 = st.columns(4)

    current_portfolio_val = portfolio_df["market_value"].sum()
    daily_pnl = portfolio_df["market_value"].sum() * 0.012  # Mock daily change

    with c1:
        st.markdown(
            ThemeEngine.render_kpi_card(
                "Total AUM", f"€{current_portfolio_val / 1e6:.2f}M", "+2.4%"
            ),
            unsafe_allow_html=True,
        )
    with c2:
        st.markdown(
            ThemeEngine.render_kpi_card("Daily P&L", f"€{daily_pnl:,.0f}", "+1.2%"),
            unsafe_allow_html=True,
        )
    with c3:
        st.markdown(
            ThemeEngine.render_kpi_card(
                "Sharpe Ratio", "1.85", "+0.05", "Risk Adj. Return"
            ),
            unsafe_allow_html=True,
        )
    with c4:
        st.markdown(
            ThemeEngine.render_kpi_card(
                "Active Alerts", str(len(st.session_state.notifications)), None
            ),
            unsafe_allow_html=True,
        )

    st.markdown("---")

    # 2. Main Visuals
    col_main, col_side = st.columns([2, 1])

    with col_main:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### 💹 Market Performance (Real-time)")

        # Animated Plotly Chart
        fig = go.Figure()

        # Determine sampling for animation frames to keep performance high
        anim_df = market_df.tail(200)

        fig.add_trace(
            go.Scatter(
                x=anim_df["date"],
                y=anim_df["price"],
                mode="lines",
                fill="tozeroy",
                line=dict(color="#00d4ff", width=2),
                name="Market Index",
            )
        )

        fig.update_layout(
            template="plotly_dark",
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)",
            margin=dict(l=0, r=0, t=10, b=0),
            height=350,
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=True, gridcolor="rgba(255,255,255,0.1)"),
        )
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
        st.markdown("</div>", unsafe_allow_html=True)

        # Secondary Row: Volume & RSI
        c_sub1, c_sub2 = st.columns(2)
        with c_sub1:
            st.markdown(
                "<div class='glass-card' style='margin-top:20px'>",
                unsafe_allow_html=True,
            )
            st.markdown("#### Trading Volume")
            fig_vol = px.bar(
                market_df.tail(50),
                x="date",
                y="volume",
                color="volume",
                color_continuous_scale="Bluyl",
            )
            fig_vol.update_layout(
                template="plotly_dark",
                height=200,
                margin=dict(l=0, r=0, t=10, b=0),
                showlegend=False,
                paper_bgcolor="rgba(0,0,0,0)",
                plot_bgcolor="rgba(0,0,0,0)",
            )
            st.plotly_chart(
                fig_vol, use_container_width=True, config={"displayModeBar": False}
            )
            st.markdown("</div>", unsafe_allow_html=True)

        with c_sub2:
            st.markdown(
                "<div class='glass-card' style='margin-top:20px'>",
                unsafe_allow_html=True,
            )
            st.markdown("#### RSI (14)")
            rsi_df = market_df.tail(50)
            fig_rsi = go.Figure(
                go.Scatter(
                    x=rsi_df["date"],
                    y=rsi_df["rsi"],
                    line=dict(color="#7b2ff7", width=2),
                )
            )
            fig_rsi.add_hline(
                y=70, line_dash="dot", line_color="red", annotation_text="Overbought"
            )
            fig_rsi.add_hline(
                y=30, line_dash="dot", line_color="green", annotation_text="Oversold"
            )
            fig_rsi.update_layout(
                template="plotly_dark",
                height=200,
                margin=dict(l=0, r=0, t=10, b=0),
                paper_bgcolor="rgba(0,0,0,0)",
                plot_bgcolor="rgba(0,0,0,0)",
            )
            st.plotly_chart(
                fig_rsi, use_container_width=True, config={"displayModeBar": False}
            )
            st.markdown("</div>", unsafe_allow_html=True)

    with col_side:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### 📦 Sector Allocation")

        fig_pie = px.sunburst(
            portfolio_df,
            path=["sector", "ticker"],
            values="market_value",
            color="unrealized_pl_pct",
            color_continuous_scale="RdYlGn",
        )
        fig_pie.update_layout(
            template="plotly_dark",
            height=300,
            margin=dict(l=0, r=0, t=0, b=0),
            paper_bgcolor="rgba(0,0,0,0)",
        )
        st.plotly_chart(fig_pie, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown(
            "<div class='glass-card' style='margin-top:20px'>", unsafe_allow_html=True
        )
        st.markdown("#### ⚡ AI Insights")
        insights = [
            "Tech sector showing signs of rotation.",
            "Volatility spread widening on EM bonds.",
            "RSI indicates potential entry on XOM.",
        ]
        for ins in insights:
            st.info(ins, icon="🤖")
        st.markdown("</div>", unsafe_allow_html=True)

    # Comments
    ThemeEngine.render_comment_section("dashboard_main")


# --- ANALYTICS PAGE ---
def page_analytics():
    ThemeEngine.render_header("Deep Analytics", "Quantitative Research Hub")

    df = st.session_state.market_df

    # Tabbed Interface
    tab1, tab2, tab3 = st.tabs(
        ["🔮 Monte Carlo", "🌊 Volatility Surface", "📡 FFT Spectrum"]
    )

    with tab1:
        c1, c2 = st.columns([1, 3])
        with c1:
            st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
            st.markdown("#### Simulation Params")
            sim_days = st.slider("Forecast Days", 10, 365, 30)
            sim_runs = st.slider("Simulations", 100, 5000, 500)
            st.caption("Uses Geometric Brownian Motion based on historical volatility.")

            if st.button("Run Simulation", type="primary"):
                with st.spinner("Crunching numbers..."):
                    sim_data = AnalyticsLib.monte_carlo_simulation(
                        df, days=sim_days, simulations=sim_runs
                    )

                    # Plotting
                    fig_mc = go.Figure()
                    # Plot first 100 paths
                    for col in sim_data.columns[:100]:
                        fig_mc.add_trace(
                            go.Scatter(
                                y=sim_data[col],
                                mode="lines",
                                line=dict(color="rgba(0, 212, 255, 0.05)", width=1),
                                showlegend=False,
                            )
                        )

                    # Mean path
                    fig_mc.add_trace(
                        go.Scatter(
                            y=sim_data.mean(axis=1),
                            mode="lines",
                            line=dict(color="#fff", width=3),
                            name="Mean Path",
                        )
                    )

                    fig_mc.update_layout(
                        template="plotly_dark",
                        title="Monte Carlo Outcome Paths",
                        height=500,
                        paper_bgcolor="rgba(0,0,0,0)",
                        plot_bgcolor="rgba(0,0,0,0)",
                    )
                    st.session_state["fig_mc"] = fig_mc  # Cache for display
            st.markdown("</div>", unsafe_allow_html=True)

        with c2:
            st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
            if "fig_mc" in st.session_state:
                st.plotly_chart(st.session_state["fig_mc"], use_container_width=True)
            else:
                st.info("Set parameters and click 'Run Simulation'")
            st.markdown("</div>", unsafe_allow_html=True)

    with tab2:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### Volatility Cones Analysis")
        cones_df = AnalyticsLib.calculate_volatility_cones(df)

        # Visualizing Cones
        fig_cones = go.Figure()
        windows = cones_df.index

        fig_cones.add_trace(
            go.Scatter(
                x=windows,
                y=cones_df["max"],
                name="Max Vol",
                line=dict(color="red", dash="dash"),
            )
        )
        fig_cones.add_trace(
            go.Scatter(
                x=windows,
                y=cones_df["p75"],
                name="75th Percentile",
                line=dict(color="orange"),
            )
        )
        fig_cones.add_trace(
            go.Scatter(
                x=windows,
                y=cones_df["median"],
                name="Median Vol",
                line=dict(color="white", width=3),
            )
        )
        fig_cones.add_trace(
            go.Scatter(
                x=windows,
                y=cones_df["p25"],
                name="25th Percentile",
                line=dict(color="cyan"),
            )
        )
        fig_cones.add_trace(
            go.Scatter(
                x=windows,
                y=cones_df["min"],
                name="Min Vol",
                line=dict(color="green", dash="dash"),
            )
        )

        # Current realization
        fig_cones.add_trace(
            go.Scatter(
                x=windows,
                y=cones_df["current"],
                name="Current Realized",
                mode="markers+lines",
                marker=dict(size=10, color="yellow"),
            )
        )

        fig_cones.update_layout(
            template="plotly_dark",
            height=500,
            title="Realized Volatility Term Structure",
            xaxis_title="Lookback Window (Days)",
            yaxis_title="Annualized Volatility",
        )
        st.plotly_chart(fig_cones, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    with tab3:
        c_fft, c_desc = st.columns([3, 1])
        with c_fft:
            st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
            st.markdown("#### Fast Fourier Transform (Cycle Detection)")
            fft_df = AnalyticsLib.perform_fft_analysis(df)

            fig_fft = px.bar(
                fft_df.head(20),
                x="Period (Days)",
                y="Power",
                color="Power",
                color_continuous_scale="Viridis",
            )
            fig_fft.update_layout(template="plotly_dark", height=450)
            st.plotly_chart(fig_fft, use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
        with c_desc:
            st.write(
                "High power at specific days suggests recurring patterns (e.g., quarterly earnings, weekly cycles)."
            )
            st.markdown("</div>", unsafe_allow_html=True)

    # Comments
    ThemeEngine.render_comment_section("analytics_main")


# --- TRADING PAGE ---
def page_trading():
    ThemeEngine.render_header("Portfolio & Trading", "Execution Terminal")
    portfolio = st.session_state.portfolio_df

    col_port, col_exec = st.columns([2, 1])

    with col_port:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### 📂 Current Holdings")

        # Enhanced Dataframe
        st.dataframe(
            portfolio.style.format(
                {
                    "market_value": "€{:,.2f}",
                    "unrealized_pl": "€{:,.2f}",
                    "unrealized_pl_pct": "{:.2f}%",
                    "daily_change_pct": "{:.2%}",
                }
            ).background_gradient(
                subset=["unrealized_pl_pct"], cmap="RdYlGn", vmin=-10, vmax=10
            ),
            use_container_width=True,
            height=300,
        )
        st.markdown("</div>", unsafe_allow_html=True)

        # Risk Bubble Chart
        st.markdown(
            "<div class='glass-card' style='margin-top:20px'>", unsafe_allow_html=True
        )
        st.markdown("#### 🎈 Risk vs Reward Map")
        fig_bub = px.scatter(
            portfolio,
            x="daily_change_pct",
            y="unrealized_pl_pct",
            size="market_value",
            color="sector",
            hover_name="name",
            text="ticker",
            color_discrete_sequence=px.colors.qualitative.Pastel,
        )
        fig_bub.update_layout(
            template="plotly_dark",
            height=400,
            xaxis_title="Daily Volatility",
            yaxis_title="Total Return %",
        )
        st.plotly_chart(fig_bub, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    with col_exec:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### ⚡ Quick Trade")

        ticker = st.selectbox(
            "Asset", portfolio["ticker"].tolist() + ["TSLA", "AMZN", "MSFT"]
        )
        order_type = st.radio(
            "Order Type", ["Market", "Limit", "Stop Loss"], horizontal=True
        )
        side = st.selectbox("Side", ["BUY", "SELL"])

        qty = st.number_input("Quantity", 1, 10000, 10)

        price = 150.00  # Mock price
        if order_type == "Limit":
            price = st.number_input("Limit Price", 0.0, 5000.0, 150.0)

        est_total = qty * price

        st.markdown("---")
        st.markdown(
            f"<div style='display:flex; justify-content:space-between'><span>Est. Price:</span> <b>€{price:.2f}</b></div>",
            unsafe_allow_html=True,
        )
        st.markdown(
            f"<div style='display:flex; justify-content:space-between'><span>Fees (0.1%):</span> <b>€{est_total * 0.001:.2f}</b></div>",
            unsafe_allow_html=True,
        )
        st.markdown(
            f"<div style='font-size:1.5rem; color:var(--primary-color); font-weight:700; text-align:right; margin-top:10px'>€{est_total * 1.001:,.2f}</div>",
            unsafe_allow_html=True,
        )

        btn_col = "var(--success)" if side == "BUY" else "var(--danger)"

        if st.button(f"SUBMIT {side} ORDER", use_container_width=True):
            with st.spinner("Routing to exchange..."):
                time.sleep(0.5)
            st.success(f"Order Filled: {side} {qty} {ticker} @ {price}")
            st.session_state.activity_log.append(
                {
                    "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "user": st.session_state.user["name"],
                    "action": "Trade Executed",
                    "details": f"{side} {qty} {ticker} @ {price}",
                }
            )

        st.markdown("</div>", unsafe_allow_html=True)

        # Simulated Order Book
        st.markdown(
            "<div class='glass-card' style='margin-top:20px'>", unsafe_allow_html=True
        )
        st.markdown("#### 📊 Order Book Depth")

        depth_data = {
            "price": np.concatenate(
                [np.linspace(148, 149.9, 10), np.linspace(150.1, 152, 10)]
            ),
            "size": np.random.randint(100, 5000, 20),
            "side": ["Bid"] * 10 + ["Ask"] * 10,
        }
        df_depth = pd.DataFrame(depth_data)

        fig_ob = px.bar(
            df_depth,
            x="size",
            y="price",
            color="side",
            orientation="h",
            color_discrete_map={"Bid": "#00e396", "Ask": "#ff0055"},
        )
        fig_ob.update_layout(
            template="plotly_dark",
            height=300,
            margin=dict(l=0, r=0, t=0, b=0),
            bargap=0.1,
        )
        st.plotly_chart(
            fig_ob, use_container_width=True, config={"displayModeBar": False}
        )
        st.markdown("</div>", unsafe_allow_html=True)


# --- SALES & CUSTOMERS PAGE ---
def page_sales():
    ThemeEngine.render_header("Revenue Intelligence", "Cohort & Sales Analysis")
    df = st.session_state.sales_df

    # 1. Geographic & Channel Analysis
    c1, c2 = st.columns([2, 1])
    with c1:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### 🌍 Global Revenue Map")
        # Aggregation
        geo_agg = df.groupby("region")["amount"].sum().reset_index()
        fig_map = px.choropleth(
            geo_agg,
            locations="region",
            locationmode="country names",
            color="amount",
            color_continuous_scale="Plasma",
            scope="world",
        )
        # Since 'region' are broad names, let's use a simpler bar chart for robustness if geo fails matching
        fig_map = px.bar(
            geo_agg,
            x="region",
            y="amount",
            color="amount",
            color_continuous_scale="Plasma",
        )
        fig_map.update_layout(template="plotly_dark", height=350)
        st.plotly_chart(fig_map, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    with c2:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### 📢 Channel Mix")
        chan_agg = df.groupby("channel")["amount"].sum().reset_index()
        fig_pie = px.doughnut(chan_agg, values="amount", names="channel", hole=0.6)
        fig_pie.update_layout(
            template="plotly_dark",
            height=350,
            showlegend=True,
            legend=dict(orientation="h", y=-0.1),
        )
        st.plotly_chart(fig_pie, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("---")

    # 2. Cohort Analysis (The Complex Bit)
    st.markdown("### 👥 Retention Cohorts")
    st.markdown("<div class='glass-card'>", unsafe_allow_html=True)

    # Pivot table for cohorts
    cohort_counts = (
        df.groupby(["cohort_month", "month"])["customer_id"].nunique().reset_index()
    )
    cohort_counts["period_number"] = (
        cohort_counts.month - cohort_counts.cohort_month
    ).apply(lambda x: x.n)

    cohort_pivot = cohort_counts.pivot_table(
        index="cohort_month", columns="period_number", values="customer_id"
    )
    cohort_size = cohort_pivot.iloc[:, 0]
    retention_matrix = cohort_pivot.divide(cohort_size, axis=0)

    # Heatmap
    fig_coh = px.imshow(
        retention_matrix,
        text_auto=".0%",
        color_continuous_scale="Blues",
        labels=dict(x="Months Since Acquisition", y="Cohort", color="Retention"),
    )
    fig_coh.update_layout(template="plotly_dark", height=500)
    st.plotly_chart(fig_coh, use_container_width=True)
    st.markdown("</div>", unsafe_allow_html=True)


# --- REPORT BUILDER (PYQT5) ---
def page_reports():
    ThemeEngine.render_header("Report Generator", "High-Fidelity PDF Engine")

    c_config, c_preview = st.columns([1, 2])

    with c_config:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### 📄 Configuration")

        rep_title = st.text_input("Report Title", "Monthly Investment Memorandum")
        rep_author = st.text_input("Author", st.session_state.user["name"])

        st.markdown("##### Sections to Include")
        inc_summary = st.checkbox("Executive Summary", True)
        inc_kpis = st.checkbox("Key Performance Indicators", True)
        inc_holdings = st.checkbox("Top Holdings Table", True)
        inc_charts = st.checkbox("Performance Charts", True)
        inc_disclaimer = st.checkbox("Legal Disclaimer", True)

        st.markdown("---")
        notes = st.text_area(
            "Analyst Commentary",
            "The portfolio outperformed the benchmark by 150bps due to strong selection in the tech sector...",
        )

        generate_btn = st.button(
            "Generate PDF", type="primary", use_container_width=True
        )
        st.markdown("</div>", unsafe_allow_html=True)

    with c_preview:
        if generate_btn:
            with st.spinner("Initializing PyQt5 Subprocess..."):
                # 1. Prepare Data for Report
                port = st.session_state.portfolio_df
                total_aum = port["market_value"].sum()
                top_holdings = port.nlargest(5, "market_value")[
                    ["name", "ticker", "sector", "market_value", "unrealized_pl_pct"]
                ]

                # 2. Construct HTML Template (CSS for Print)
                html_template = f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <style>
                        body {{ font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif; padding: 40px; color: #333; line-height: 1.6; }}
                        .header {{ text-align: center; border-bottom: 2px solid #00d4ff; padding-bottom: 20px; margin-bottom: 30px; }}
                        .header h1 {{ margin: 0; color: #2c3e50; font-size: 32px; }}
                        .header p {{ color: #7f8c8d; margin: 5px 0 0 0; }}
                        
                        .section {{ margin-bottom: 30px; }}
                        .section h2 {{ color: #00d4ff; border-bottom: 1px solid #eee; padding-bottom: 10px; font-size: 20px; }}
                        
                        .kpi-container {{ display: flex; justify-content: space-between; gap: 20px; margin-bottom: 30px; }}
                        .kpi-box {{ flex: 1; background: #f8f9fa; padding: 15px; border-radius: 8px; text-align: center; border: 1px solid #e9ecef; }}
                        .kpi-val {{ font-size: 24px; font-weight: bold; color: #2c3e50; }}
                        .kpi-lbl {{ font-size: 12px; color: #7f8c8d; text-transform: uppercase; letter-spacing: 1px; }}
                        
                        table {{ width: 100%; border-collapse: collapse; margin-top: 15px; font-size: 12px; }}
                        th {{ background: #f1f3f5; text-align: left; padding: 10px; border-bottom: 2px solid #dee2e6; color: #495057; }}
                        td {{ padding: 10px; border-bottom: 1px solid #dee2e6; }}
                        tr:nth-child(even) {{ background-color: #f8f9fa; }}
                        
                        .footer {{ margin-top: 50px; font-size: 10px; color: #adb5bd; text-align: center; border-top: 1px solid #eee; padding-top: 20px; }}
                        .commentary {{ background: #fff3cd; padding: 15px; border-radius: 5px; border-left: 4px solid #ffc107; font-style: italic; }}
                    </style>
                </head>
                <body>
                    <div class="header">
                        <h1>{rep_title}</h1>
                        <p>Generated by Aurora Engine | {datetime.now().strftime("%B %d, %Y")}</p>
                        <p>Prepared by: {rep_author}</p>
                    </div>
                """

                if inc_kpis:
                    html_template += f"""
                    <div class="kpi-container">
                        <div class="kpi-box">
                            <div class="kpi-lbl">Total AUM</div>
                            <div class="kpi-val">€{total_aum:,.0f}</div>
                        </div>
                        <div class="kpi-box">
                            <div class="kpi-lbl">Holdings</div>
                            <div class="kpi-val">{len(port)}</div>
                        </div>
                        <div class="kpi-box">
                            <div class="kpi-lbl">Avg Return</div>
                            <div class="kpi-val">{port["unrealized_pl_pct"].mean():.2f}%</div>
                        </div>
                    </div>
                    """

                if inc_summary:
                    html_template += f"""
                    <div class="section">
                        <h2>Analyst Commentary</h2>
                        <div class="commentary">{notes}</div>
                    </div>
                    """

                if inc_holdings:
                    # Convert dataframe to HTML table
                    table_html = top_holdings.to_html(index=False, classes="", border=0)
                    html_template += f"""
                    <div class="section">
                        <h2>Top Holdings Breakdown</h2>
                        {table_html}
                    </div>
                    """

                # 2.1 Add Custom Charts from Visual Builder
                if (
                    inc_charts
                    and "my_charts" in st.session_state
                    and st.session_state.my_charts
                ):
                    st.write(
                        f"DEBUG: Adding {len(st.session_state.my_charts)} charts to PDF"
                    )  # Debug Line
                    html_template += (
                        """<div class="section"><h2>Custom Analytics</h2>"""
                    )
                    for chart in st.session_state.my_charts:
                        # Prefer local file path if available (QWebEngine handles file:// robustly)
                        img_tag = ""
                        if chart.get("image_path") and os.path.exists(
                            chart["image_path"]
                        ):
                            # use file:// URL
                            img_url = "file://" + chart["image_path"].replace("\\", "/")
                            img_tag = f'<img src="{img_url}" style="max-width:100%; height:auto; border:1px solid #ddd; border-radius:4px;"/>'
                        elif chart.get("image"):
                            # fallback to base64 inline (kept as last resort)
                            b64_chart = base64.b64encode(chart["image"]).decode("utf-8")
                            img_tag = f'<img src="data:image/png;base64,{b64_chart}" style="max-width:100%; height:auto; border:1px solid #ddd; border-radius:4px;"/>'
                        else:
                            img_tag = (
                                "<div style='color:#888'>[No image available]</div>"
                            )

                        # Wrap each chart in its own print-friendly block with a break so pages render fully
                        html_template += f"""
                        <div style="margin-bottom:20px; page-break-after:always; break-inside:avoid;">
                            <h3 style="color:#2c3e50; font-size:16px; margin-bottom:8px;">{chart.get("title", "Chart")}</h3>
                            {img_tag}
                        </div>
                        """

                if inc_disclaimer:
                    html_template += """
                    <div class="footer">
                        CONFIDENTIALITY NOTICE: The contents of this document are intended solely for the addressee. 
                        Past performance is not indicative of future results. Generated via Aurora Dashboard v8.0 Enterprise.
                    </div>
                    """

                html_template += "</body></html>"

                # 3. Call PDF Engine
                success, result = PDFEngine.generate(html_template, f"{rep_title}.pdf")

                if success:
                    st.success("PDF Report Successfully Generated!")
                    st.markdown(f"**Size:** {len(result) / 1024:.1f} KB")

                    # Actions Row
                    ac1, ac2 = st.columns(2)
                    with ac1:
                        st.download_button(
                            label="⬇ Download Final PDF",
                            data=result,
                            file_name=f"{rep_title.replace(' ', '_')}.pdf",
                            mime="application/pdf",
                            type="primary",
                            use_container_width=True,
                        )
                    with ac2:
                        if st.button(
                            "📅 Schedule this Report", use_container_width=True
                        ):
                            st.session_state.scheduler_defaults = {
                                "name": rep_title,
                                "format": "PDF Report",
                                "source": "Report Generator",
                            }
                            st.success(
                                "Configuration saved! Redirecting to Scheduler..."
                            )
                            time.sleep(1)
                            st.switch_page(
                                "page_scheduler"
                            )  # Note: using page_scheduler won't work with single script router.
                            # We need to set 'selected_module' but our router uses a radio.
                            # Workaround: Set a flag and rerun, or just instruct user.
                            # Better: We'll modify the router to respect a session override,
                            # but for now let's just use st.rerun() and hope user navigates manually?
                            # Wait, we can't force navigation easily in this radio setup without `st.session_state.module_selection`.
                            # Let's simple communicate.
                            st.info(
                                "Please navigate to the '📅 Scheduler' page to finalize."
                            )

                    # Preview frame (using iframe)
                    b64_pdf = base64.b64encode(result).decode("utf-8")
                    pdf_display = f'<iframe src="data:application/pdf;base64,{b64_pdf}" width="100%" height="600" type="application/pdf"></iframe>'
                    st.markdown(pdf_display, unsafe_allow_html=True)
                else:
                    st.error("PDF Generation Failed")
                    st.code(result)
                    st.warning(
                        "Check if PyQt5 and QtWebEngine are installed in the environment."
                    )
        else:
            st.info("Configure report settings and click 'Generate PDF' to preview.")

    def build_composed_pdf_from_my_charts(report_title="Composed Charts Report"):
        """
        Builds a single PDF containing all charts saved in st.session_state.my_charts.
        Returns (success: bool, pdf_bytes_or_error)
        """
        if "my_charts" not in st.session_state or not st.session_state.my_charts:
            return False, "No charts in composer."

        # Build HTML with file:// preferred images as above
        html = "<!doctype html><html><head><meta charset='utf-8'><style>body{font-family:Arial; padding:30px;} img{max-width:100%; height:auto; border:1px solid #ddd; border-radius:4px;}</style></head><body>"
        html += f"<h1>{report_title}</h1><p>Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>"
        for c in st.session_state.my_charts:
            if c.get("image_path") and os.path.exists(c["image_path"]):
                img_url = "file://" + c["image_path"].replace("\\", "/")
                html += f"<div style='margin-top:20px; page-break-after:always;'><h3>{c.get('title', 'Chart')}</h3><img src='{img_url}'/></div>"
            elif c.get("image"):
                b64 = base64.b64encode(c["image"]).decode("utf-8")
                html += f"<div style='margin-top:20px; page-break-after:always;'><h3>{c.get('title', 'Chart')}</h3><img src='data:image/png;base64,{b64}'/></div>"
        html += "</body></html>"

        return PDFEngine.generate(html, f"{report_title}.pdf")

    title = "Composed Charts Report"
    success, pdf_data_or_err = build_composed_pdf_from_my_charts(title)
    if not success:
        st.error(f"Could not build PDF: {pdf_data_or_err}")
    else:
        # Build email
        msg = EmailMessage()
        msg["Subject"] = f"{title} — Automated Report"
        msg["From"] = st.session_state.smtp_config.get("from", "reports@aurora.io")
        msg["To"] = ""  # ensure string or list
        msg.set_content("Please find the composed charts report attached.")
        msg.add_attachment(
            pdf_data_or_err,
            maintype="application",
            subtype="pdf",
            filename=f"{title}.pdf",
        )
        # send via SMTP (simple)
        try:
            s = smtplib.SMTP(
                st.session_state.smtp_config["host"],
                st.session_state.smtp_config["port"],
            )
            s.starttls()
            if st.session_state.smtp_config["user"]:
                s.login(
                    st.session_state.smtp_config["user"],
                    st.session_state.smtp_config["pass"],
                )
            s.send_message(msg)
            s.quit()
            st.success("Email sent.")
        except Exception as e:
            st.error(f"Email send failed: {e}")


# --- SETTINGS & ADMIN PAGE ---
def page_settings():
    ThemeEngine.render_header("System Settings", "Configuration Panel")

    t1, t2 = st.tabs(["🎨 Appearance", "🛡 Admin"])

    with t1:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### Theme Selection")
        current = st.session_state.theme
        options = list(ThemeEngine.THEMES.keys())
        idx = options.index(current) if current in options else 0

        new_theme = st.selectbox("Active Theme", options, index=idx)
        if new_theme != current:
            st.session_state.theme = new_theme
            st.rerun()

        st.markdown("#### Density")
        st.radio("Layout Density", ["Compact", "Comfortable"], horizontal=True)
        st.markdown("</div>", unsafe_allow_html=True)

    with t2:
        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
        st.markdown("#### Data Management")
        c1, c2 = st.columns(2)
        with c1:
            if st.button("Reset Session State", type="secondary"):
                for key in list(st.session_state.keys()):
                    del st.session_state[key]
                st.rerun()
        with c2:
            if st.button("Clear Cache", type="secondary"):
                st.cache_data.clear()
                st.success("Cache Cleared!")

        st.markdown("---")
        st.markdown("#### System Activity Log")
        log_df = pd.DataFrame(st.session_state.activity_log)
        st.dataframe(log_df, use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)


# --- DATA WRANGLING PAGE ---
def page_data_wrangling():
    ThemeEngine.render_header("Data Studio", "Wrangling & Export Hub")

    # 1. Source Selection
    st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
    st.markdown("#### 1. Select Data Source")

    source_map = {
        "Market Data (Real-time)": st.session_state.market_df,
        "Sales Transactions": st.session_state.sales_df,
        "Portfolio Holdings": st.session_state.portfolio_df,
        "Company Registry": st.session_state.company_df,
    }

    source_name = st.selectbox("Choose Dataset", list(source_map.keys()))
    df = source_map[source_name].copy()

    # File Upload Option
    uploaded_file = st.file_uploader(
        "Or upload your own (CSV/Excel)", type=["csv", "xlsx"]
    )
    if uploaded_file:
        try:
            if uploaded_file.name.endswith(".csv"):
                df = pd.read_csv(uploaded_file)
            else:
                df = pd.read_excel(uploaded_file)
            st.success(f"Loaded {len(df)} rows from {uploaded_file.name}")
        except Exception as e:
            st.error(f"Error loading file: {e}")

    st.markdown("</div>", unsafe_allow_html=True)

    # 2. Transformation
    if not df.empty:
        c1, c2 = st.columns([1, 2])

        with c1:
            st.markdown(
                "<div class='glass-card' style='height:100%'>", unsafe_allow_html=True
            )
            st.markdown("#### 2. Configure Columns")

            all_cols = df.columns.tolist()
            selected_cols = st.multiselect(
                "Select Columns to Include",
                all_cols,
                default=all_cols[: min(5, len(all_cols))],
            )

            if not selected_cols:
                st.warning("Please select at least one column.")
                selected_cols = all_cols

            st.markdown("---")
            st.markdown("#### 3. Filtering")

            filter_col = st.selectbox("Filter Column", ["None"] + all_cols)

            filtered_df = df[selected_cols].copy()

            if filter_col != "None":
                dtype = df[filter_col].dtype
                if np.issubdtype(dtype, np.number):
                    min_val = float(df[filter_col].min())
                    max_val = float(df[filter_col].max())
                    val_range = st.slider(
                        f"Range for {filter_col}", min_val, max_val, (min_val, max_val)
                    )
                    filtered_df = filtered_df[
                        (df[filter_col] >= val_range[0])
                        & (df[filter_col] <= val_range[1])
                    ]
                else:
                    unique_vals = df[filter_col].unique()
                    if len(unique_vals) < 50:
                        selected_vals = st.multiselect(
                            f"Values for {filter_col}", unique_vals, default=unique_vals
                        )
                        filtered_df = filtered_df[df[filter_col].isin(selected_vals)]
                    else:
                        text_search = st.text_input(f"Search in {filter_col}")
                        if text_search:
                            filtered_df = filtered_df[
                                df[filter_col]
                                .astype(str)
                                .str.contains(text_search, case=False)
                            ]

            st.markdown("</div>", unsafe_allow_html=True)

        with c2:
            st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
            st.markdown(f"#### 4. Data Preview ({len(filtered_df)} rows)")
            st.dataframe(filtered_df, use_container_width=True, height=400)
            st.markdown("</div>", unsafe_allow_html=True)

    # 3. Export & Action
    st.markdown(
        "<div class='glass-card' style='margin-top:20px'>", unsafe_allow_html=True
    )
    st.markdown("#### 5. Actions")

    ac1, ac2 = st.columns(2)
    with ac1:
        # Excel Export
        buffer = io.BytesIO()
        with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
            filtered_df.to_excel(writer, sheet_name="Data", index=False)

        st.download_button(
            label="⬇ Export to Excel",
            data=buffer,
            file_name=f"aurora_export_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",
            mime="application/vnd.ms-excel",
            type="primary",
        )

    with ac2:
        # Scheduler Mockup
        with st.expander("✉ Email Scheduler"):
            email_rec = st.text_input("Recipients (comma separated)")
            freq = st.selectbox(
                "Frequency", ["Daily @ 9:00 AM", "Weekly (Mon)", "Monthly (1st)"]
            )
            if st.button("Schedule Report"):
                st.session_state.activity_log.append(
                    {
                        "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "user": st.session_state.user["name"],
                        "action": "Scheduled Report Created",
                        "details": f"Dataset: {source_name} | Freq: {freq}",
                    }
                )
                st.success("Report scheduled successfully! (See Activity Log)")

    st.markdown("</div>", unsafe_allow_html=True)


# --- VISUAL BUILDER PAGE ---
# --- VISUAL BUILDER PAGE (All-in) ---
def page_visual_builder():
    """
    All-in Visual Builder:
      - Smart multi-dataset join assistant (fuzzy suggestions + preview)
      - Full filter/query builder (grouping, AND/OR)
      - Data transformations (calc cols, rename, types, binning, split/concat, pivot/unpivot)
      - Multi-aggregation (multi-group, multiple funcs per column)
      - Chart Builder with advanced features (dual-axis, combo, trendline)
      - Super Composer: reorder, text blocks, tables, page breaks, themes
      - Export: PDF (via PDFEngine.generate), PPTX (if python-pptx installed)
      - Email send (SMTP config in st.session_state['smtp_config'])
      - Caching and lazy "Build" model, no st.experimental_rerun usage
    """

    ThemeEngine.render_header("Visual Builder", "Pro Analytics Studio")

    # --------------------------
    # session defaults & source map
    # --------------------------
    st.session_state.setdefault("saved_views", {})
    st.session_state.setdefault("my_composer", [])  # more structured composer items
    st.session_state.setdefault("vb_join_steps", [])
    st.session_state.setdefault("vb_filters", [])
    st.session_state.setdefault("vb_agg_rules", [])
    st.session_state.setdefault("vb_transform_steps", [])
    st.session_state.setdefault("scheduler_defaults", {})
    st.session_state.setdefault("smtp_config", st.session_state.get("smtp_config", {}))
    st.session_state.setdefault("_vb_cached_working_df", None)
    st.session_state.setdefault("_vb_cache_stamp", None)

    source_map = {
        "Market Data": st.session_state.get("market_df", pd.DataFrame()),
        "Sales Transactions": st.session_state.get("sales_df", pd.DataFrame()),
        "Portfolio": st.session_state.get("portfolio_df", pd.DataFrame()),
        "Company Registry": st.session_state.get("company_df", pd.DataFrame()),
    }

    # --------------------------
    # helper utilities
    # --------------------------
    def safe_cols(df):
        return list(df.columns) if isinstance(df, pd.DataFrame) and not df.empty else []

    def get_similarity_pairs(left_cols, right_cols, top_n=10):
        """Return list of suggested (left, right, score) tuples using difflib.SequenceMatcher"""
        pairs = []
        for lc in left_cols:
            for rc in right_cols:
                score = difflib.SequenceMatcher(None, lc.lower(), rc.lower()).ratio()
                pairs.append((lc, rc, score))
        pairs.sort(key=lambda x: x[2], reverse=True)
        return pairs[:top_n]

    def write_temp_png_from_bytes(bts):
        if not bts:
            return None
        try:
            tmp_dir = tempfile.gettempdir()
            path = os.path.join(tmp_dir, f"vb_chart_{uuid.uuid4().hex}.png")
            with open(path, "wb") as fh:
                fh.write(bts)
            return path
        except Exception:
            return None

    def build_pdf_from_composer(
        items, title="Composed Charts Report", theme="professional"
    ):
        # Build HTML with sections; prefer file:// images
        html = [
            "<html><head><meta charset='utf-8'><style>body{font-family:Arial;padding:20px;}"
        ]
        # theme small styles
        if theme == "dark":
            html.append("body{background:#111;color:#eee} h1,h2{color:#fff}")
        html.append(
            "img{max-width:100%;height:auto;margin:10px 0;} section{margin-bottom:30px; page-break-after:always;}</style></head><body>"
        )
        html.append(
            f"<h1>{title}</h1><p>Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>"
        )
        for it in items:
            html.append("<section>")
            html.append(f"<h2>{it.get('title', 'Untitled')}</h2>")
            if it.get("subtitle"):
                html.append(f"<h4>{it.get('subtitle')}</h4>")
            if it.get("description"):
                html.append(f"<p>{it.get('description')}</p>")
            # table block support
            if it.get("type") == "table" and isinstance(
                it.get("table_df"), pd.DataFrame
            ):
                # render table as HTML
                table_html = it["table_df"].to_html(
                    index=False, classes="vb-table", border=0
                )
                html.append(table_html)
            # image / chart
            if it.get("type") in ("chart", "image"):
                if it.get("image_path") and os.path.exists(it["image_path"]):
                    src = "file://" + it["image_path"].replace("\\", "/")
                    html.append(f"<img src='{src}' alt='chart'/>")
                elif it.get("image"):
                    b64 = base64.b64encode(it["image"]).decode("utf-8")
                    html.append(f"<img src='data:image/png;base64,{b64}' alt='chart'/>")
                else:
                    html.append("<p><i>[No image available]</i></p>")
            # page break control
            if it.get("page_break", True):
                html.append("<div style='page-break-after:always'></div>")
            html.append("</section>")
        html.append("</body></html>")
        final_html = "".join(html)
        return PDFEngine.generate(final_html, f"composed_{uuid.uuid4().hex}.pdf")

    def build_pptx_from_composer(items, title="Composed Charts Report"):
        if not HAS_PPTX:
            return False, "python-pptx not installed"
        prs = Presentation()
        # title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[
            1
        ].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        for it in items:
            # simple slide per item
            slide_layout = prs.slide_layouts[5]  # blank
            slide = prs.slides.add_slide(slide_layout)
            left = Inches(0.5)
            top = Inches(0.4)
            width = Inches(9)
            # add title
            tx = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            tx.text = it.get("title", "")
            # image if present
            if it.get("image_path") and os.path.exists(it["image_path"]):
                slide.shapes.add_picture(
                    it["image_path"], Inches(0.5), Inches(1.2), width=Inches(9)
                )
            elif it.get("image"):
                # write temp and add
                temp_img = write_temp_png_from_bytes(it["image"])
                if temp_img:
                    slide.shapes.add_picture(
                        temp_img, Inches(0.5), Inches(1.2), width=Inches(9)
                    )
        # export pptx bytes
        out = io.BytesIO()
        prs.save(out)
        return True, out.getvalue()

    # --------------------------
    # Layout: left pane controls, big canvas right
    # --------------------------
    # Use columns to get a two-pane feel
    left_col, right_col = st.columns([1, 2])

    # --------------------------
    # LEFT: Data Source, Joins, Smart Assistant
    # --------------------------
    with left_col:
        st.markdown("## Data & Joins")
        ds_keys = list(source_map.keys())
        primary_ds = st.selectbox("Primary dataset", ds_keys, index=0)
        primary_df = source_map.get(primary_ds, pd.DataFrame()).copy()
        st.caption(
            f"{primary_ds}: {len(primary_df)} rows • {len(primary_df.columns)} cols"
        )

        # Smart suggestions
        st.markdown("### Smart Join Assistant")
        other_ds_for_suggest = st.selectbox(
            "Pick another dataset to suggest keys",
            [d for d in ds_keys if d != primary_ds],
        )
        other_df = source_map.get(other_ds_for_suggest, pd.DataFrame())
        left_cols = safe_cols(primary_df)
        right_cols = safe_cols(other_df)
        suggestions = get_similarity_pairs(left_cols, right_cols, top_n=12)
        if suggestions:
            st.markdown("Top suggested join keys (similarity score 0..1):")
            for a, b, s in suggestions[:8]:
                st.markdown(f"- `{a}` ⟷ `{b}` — **{s:.2f}**")
        else:
            st.info("No candidate suggestions (one of the datasets is empty).")

        st.markdown("#### Build join chain (ordered)")
        # show chain
        if st.session_state["vb_join_steps"]:
            for i, stp in enumerate(st.session_state["vb_join_steps"]):
                st.markdown(
                    f"{i + 1}. {primary_ds} ⟷ {stp['other_ds']} on `{stp['left_on']}` = `{stp['right_on']}` ({stp['how']})"
                )
                cols = st.columns([1, 1, 1])
                if cols[0].button("Up", key=f"join_up_{i}") and i > 0:
                    (
                        st.session_state["vb_join_steps"][i - 1],
                        st.session_state["vb_join_steps"][i],
                    ) = (
                        st.session_state["vb_join_steps"][i],
                        st.session_state["vb_join_steps"][i - 1],
                    )
                    st.success("Moved up")
                if (
                    cols[1].button("Down", key=f"join_down_{i}")
                    and i < len(st.session_state["vb_join_steps"]) - 1
                ):
                    (
                        st.session_state["vb_join_steps"][i],
                        st.session_state["vb_join_steps"][i + 1],
                    ) = (
                        st.session_state["vb_join_steps"][i + 1],
                        st.session_state["vb_join_steps"][i],
                    )
                    st.success("Moved down")
                if cols[2].button("Remove", key=f"join_rm_{i}"):
                    st.session_state["vb_join_steps"].pop(i)
                    st.success("Removed step")

        # form to add join
        with st.form("vb_add_join_form", clear_on_submit=True):
            other_ds = st.selectbox(
                "Join with dataset",
                [d for d in ds_keys if d != primary_ds],
                key="vb_join_other",
            )
            right_df = source_map.get(other_ds, pd.DataFrame())
            left_key = st.selectbox(
                "Primary key (left)",
                ["--select--"] + safe_cols(primary_df),
                key="vb_join_left",
            )
            right_key = st.selectbox(
                "Other key (right)",
                ["--select--"] + safe_cols(right_df),
                key="vb_join_right",
            )
            how = st.selectbox(
                "Join type", ["inner", "left", "right", "outer"], index=0
            )
            add_join = st.form_submit_button("Add Join Step")
            if add_join:
                if left_key == "--select--" or right_key == "--select--":
                    st.error("Select valid keys.")
                else:
                    st.session_state["vb_join_steps"].append(
                        {
                            "other_ds": other_ds,
                            "left_on": left_key,
                            "right_on": right_key,
                            "how": how,
                        }
                    )
                    st.success("Join step added to chain.")

        # preview join result (build working_df but don't cache until user Build)
        st.markdown("**Quick preview of joined data (top 10 rows)**")
        working_df_preview = primary_df.copy()
        jerr = []
        for sidx, sstep in enumerate(st.session_state["vb_join_steps"]):
            right_df = source_map.get(sstep["other_ds"], pd.DataFrame()).copy()
            try:
                working_df_preview = working_df_preview.merge(
                    right_df,
                    left_on=sstep["left_on"],
                    right_on=sstep["right_on"],
                    how=sstep["how"],
                    suffixes=("", f"_{sstep['other_ds'].replace(' ', '_')}"),
                )
            except Exception as e:
                jerr.append(f"Step {sidx + 1}: {e}")
        if jerr:
            for er in jerr:
                st.error(er)
        st.dataframe(working_df_preview.head(10), use_container_width=True)

        st.markdown("---")
        # Build and cache working df
        if st.button("Build working dataset (apply join chain)"):
            try:
                working_df = primary_df.copy()
                for sstep in st.session_state["vb_join_steps"]:
                    right_df = source_map.get(sstep["other_ds"], pd.DataFrame()).copy()
                    working_df = working_df.merge(
                        right_df,
                        left_on=sstep["left_on"],
                        right_on=sstep["right_on"],
                        how=sstep["how"],
                        suffixes=("", f"_{sstep['other_ds'].replace(' ', '_')}"),
                    )
                st.session_state["_vb_cached_working_df"] = working_df
                st.session_state["_vb_cache_stamp"] = datetime.now().isoformat()
                st.success("Working dataset built and cached.")
            except Exception as e:
                st.error(f"Join/build failed: {e}")

    # --------------------------
    # RIGHT: Canvas (tabs for refinement, transform, agg, charts, composer)
    # --------------------------
    with right_col:
        tabs = st.tabs(
            [
                "Refine",
                "Transform",
                "Aggregate",
                "Smart Templates",
                "Charts",
                "Composer",
            ]
        )
        # load working df from cache if available else primary preview
        cached = st.session_state.get("_vb_cached_working_df")
        if cached is None:
            working_df = (
                working_df_preview.copy()
                if "working_df_preview" in locals()
                else primary_df.copy()
            )
        else:
            working_df = cached.copy()

        # --------------------------
        # Refine Tab - advanced filters & templates
        # --------------------------
        with tabs[0]:
            st.markdown("## Refinement & Filters")
            st.caption(
                f"Working rows: {len(working_df)} • cols: {len(working_df.columns)}"
            )
            # select columns
            all_cols = safe_cols(working_df)
            sel_cols = st.multiselect("Columns to keep", all_cols, default=all_cols)
            refined = working_df.copy()
            if sel_cols:
                refined = refined[sel_cols].copy()

            # filter templates
            st.markdown("### Filter Templates")
            if st.button("Exclude nulls (all columns)"):
                refined = refined.dropna()
                st.success("Dropped rows containing nulls.")
            if st.button("Top 100 by a numeric column"):
                num_cols = [
                    c for c in all_cols if pd.api.types.is_numeric_dtype(working_df[c])
                ]
                if num_cols:
                    top_col = st.selectbox(
                        "Choose column for TopN", num_cols, key="topn_col"
                    )
                    try:
                        refined = refined.sort_values(top_col, ascending=False).head(
                            100
                        )
                        st.success("Filtered top 100 rows.")
                    except Exception as e:
                        st.error(f"TopN failed: {e}")
                else:
                    st.info("No numeric columns to apply TopN.")

            st.markdown("### Query Builder (Grouped conditions)")
            st.write(
                "Use the form to add rules. You can combine later in AND / OR groups."
            )
            if "vb_query_groups" not in st.session_state:
                st.session_state[
                    "vb_query_groups"
                ] = []  # list of {op:'AND'/'OR', rules:[{col,op,val}]}

            # show groups
            for gi, grp in enumerate(st.session_state["vb_query_groups"]):
                st.markdown(
                    f"Group {gi + 1} — {grp.get('op', 'AND')} — {len(grp.get('rules', []))} rules"
                )
                if st.button("Remove group", key=f"rm_grp_{gi}"):
                    st.session_state["vb_query_groups"].pop(gi)
                    st.success("Group removed.")

            with st.form("add_group_form", clear_on_submit=True):
                gop = st.selectbox("Group operator", ["AND", "OR"], index=0)
                # build rule list in this form
                col = st.selectbox("Column", ["--select--"] + all_cols, key="grp_col")
                if col != "--select--":
                    dtype = refined[col].dtype
                    if np.issubdtype(dtype, np.number):
                        rop = st.selectbox(
                            "Operator",
                            ["==", "!=", ">", "<", ">=", "<=", "between"],
                            key="grp_num_op",
                        )
                        if rop == "between":
                            lo = float(refined[col].min())
                            hi = float(refined[col].max())
                            rval = st.slider(
                                "Range", lo, hi, (lo, hi), key="grp_num_rng"
                            )
                        else:
                            rval = st.number_input(
                                "Value",
                                value=float(
                                    refined[col].median()
                                    if not refined[col].isna().all()
                                    else 0.0
                                ),
                                key="grp_num_val",
                            )
                    else:
                        rop = st.selectbox(
                            "Operator", ["contains", "==", "!="], key="grp_str_op"
                        )
                        rval = st.text_input("Value", key="grp_str_val")
                else:
                    rop, rval = None, None
                submit_group = st.form_submit_button("Add Group")
                if submit_group:
                    if col == "--select--":
                        st.error("Pick a column.")
                    else:
                        st.session_state["vb_query_groups"].append(
                            {"op": gop, "rules": [{"col": col, "op": rop, "val": rval}]}
                        )
                        st.success("Added group")

            # apply query groups
            try:
                if st.session_state["vb_query_groups"]:
                    combined_mask = None
                    for grp in st.session_state["vb_query_groups"]:
                        grp_mask = None
                        for r in grp["rules"]:
                            if r["op"] == "between" and isinstance(
                                r["val"], (list, tuple)
                            ):
                                m = (refined[r["col"]] >= r["val"][0]) & (
                                    refined[r["col"]] <= r["val"][1]
                                )
                            elif r["op"] == "==":
                                m = refined[r["col"]] == r["val"]
                            elif r["op"] == "!=":
                                m = refined[r["col"]] != r["val"]
                            elif r["op"] == ">":
                                m = refined[r["col"]] > r["val"]
                            elif r["op"] == "<":
                                m = refined[r["col"]] < r["val"]
                            elif r["op"] == ">=":
                                m = refined[r["col"]] >= r["val"]
                            elif r["op"] == "<=":
                                m = refined[r["col"]] <= r["val"]
                            elif r["op"] == "contains":
                                m = (
                                    refined[r["col"]]
                                    .astype(str)
                                    .str.contains(str(r["val"]), case=False, na=False)
                                )
                            else:
                                m = pd.Series([True] * len(refined))
                            if grp_mask is None:
                                grp_mask = m
                            else:
                                grp_mask = grp_mask & m
                        if combined_mask is None:
                            combined_mask = grp_mask
                        else:
                            if grp["op"] == "AND":
                                combined_mask = combined_mask & grp_mask
                            else:
                                combined_mask = combined_mask | grp_mask
                    refined = refined[combined_mask]
            except Exception as e:
                st.error(f"Query apply failed: {e}")

            st.markdown("**Refined preview**")
            st.dataframe(refined.head(50), use_container_width=True)
            # stash refined to session for later tabs
            st.session_state["_vb_refined_df"] = refined

        # --------------------------
        # Transform Tab - calculated columns, rename, pivot/unpivot
        # --------------------------
        with tabs[1]:
            st.markdown("## Transformations")
            base_df = st.session_state.get("_vb_refined_df", working_df.copy())

            st.markdown("### Calculated Column")
            with st.form("calc_col_form", clear_on_submit=True):
                new_col = st.text_input("New column name")
                expr = st.text_area(
                    "Expression (use pandas syntax, `df` is available, e.g. df['Sales'] - df['Cost'])"
                )
                add_calc = st.form_submit_button("Add Calculated Column")
                if add_calc:
                    try:
                        # create a safe local df copy
                        _tmp = base_df.copy()
                        # evaluate expression in controlled namespace
                        localns = {"df": _tmp, "np": np, "pd": pd, "math": math}
                        _tmp[new_col] = eval(expr, {}, localns)
                        base_df = _tmp
                        st.success(f"Calculated column '{new_col}' added.")
                    except Exception as e:
                        st.error(f"Calc failed: {e}")

            st.markdown("### Column Operations (rename / type / split / combine / bin)")
            # rename
            with st.form("rename_form", clear_on_submit=True):
                cols = safe_cols(base_df)
                old = st.selectbox("Column to rename", ["--select--"] + cols)
                new = st.text_input("New name")
                rename_sub = st.form_submit_button("Rename")
                if rename_sub:
                    if old == "--select--" or not new:
                        st.error("Choose and provide a new name.")
                    else:
                        base_df = base_df.rename(columns={old: new})
                        st.success("Column renamed.")

            # change type
            with st.form("type_form", clear_on_submit=True):
                cols = safe_cols(base_df)
                colt = st.selectbox(
                    "Column to change type", ["--select--"] + cols, key="tf_col"
                )
                tgt = st.selectbox(
                    "Target type", ["int", "float", "str", "datetime"], key="tf_type"
                )
                tf_sub = st.form_submit_button("Change Type")
                if tf_sub:
                    try:
                        if tgt == "int":
                            base_df[colt] = base_df[colt].astype("Int64")
                        elif tgt == "float":
                            base_df[colt] = base_df[colt].astype(float)
                        elif tgt == "str":
                            base_df[colt] = base_df[colt].astype(str)
                        elif tgt == "datetime":
                            base_df[colt] = pd.to_datetime(
                                base_df[colt], errors="coerce"
                            )
                        st.success("Type changed.")
                    except Exception as e:
                        st.error(f"Type change failed: {e}")

            # binning
            with st.form("bin_form", clear_on_submit=True):
                num_cols = [
                    c
                    for c in safe_cols(base_df)
                    if pd.api.types.is_numeric_dtype(base_df[c])
                ]
                col_bin = st.selectbox(
                    "Numeric column to bin", ["--select--"] + num_cols
                )
                bins = st.number_input(
                    "Number of bins", min_value=2, max_value=20, value=5
                )
                bcol = st.text_input("New binned column name")
                bin_sub = st.form_submit_button("Create bins")
                if bin_sub:
                    try:
                        base_df[bcol or f"{col_bin}_bin"] = pd.cut(
                            base_df[col_bin], bins=bins, labels=False
                        )
                        st.success("Binned column created.")
                    except Exception as e:
                        st.error(f"Binning failed: {e}")

            # pivot/unpivot
            with st.form("pivot_form", clear_on_submit=True):
                st.write("Pivot / Unpivot")
                p_type = st.selectbox("Action", ["pivot", "unpivot"])
                if p_type == "pivot":
                    index_cols = st.multiselect("Index (rows)", safe_cols(base_df))
                    cols_col = st.selectbox(
                        "Columns column", ["--select--"] + safe_cols(base_df)
                    )
                    vals_col = st.selectbox(
                        "Values column", ["--select--"] + safe_cols(base_df)
                    )
                    p_sub = st.form_submit_button("Pivot")
                    if p_sub:
                        try:
                            base_df = base_df.pivot_table(
                                index=index_cols,
                                columns=cols_col,
                                values=vals_col,
                                aggfunc="first",
                            ).reset_index()
                            st.success("Pivot applied.")
                        except Exception as e:
                            st.error(f"Pivot failed: {e}")
                else:
                    id_vars = st.multiselect("ID vars", safe_cols(base_df))
                    value_vars = st.multiselect(
                        "Value vars",
                        [c for c in safe_cols(base_df) if c not in id_vars],
                    )
                    u_sub = st.form_submit_button("Unpivot")
                    if u_sub:
                        try:
                            base_df = base_df.melt(
                                id_vars=id_vars,
                                value_vars=value_vars,
                                var_name="variable",
                                value_name="value",
                            )
                            st.success("Unpivot applied.")
                        except Exception as e:
                            st.error(f"Unpivot failed: {e}")

            st.markdown("**Transformed preview**")
            st.dataframe(base_df.head(50), use_container_width=True)
            st.session_state["_vb_transformed_df"] = base_df

        # --------------------------
        # Aggregate Tab
        # --------------------------
        with tabs[2]:
            st.markdown("## Aggregation Builder")
            # pick source df (transformed > refined > working)
            df_for_agg = st.session_state.get(
                "_vb_transformed_df",
                st.session_state.get("_vb_refined_df", working_df.copy()),
            )
            st.caption(
                f"Source rows: {len(df_for_agg)} • cols: {len(df_for_agg.columns)}"
            )
            cols_avail = safe_cols(df_for_agg)
            grp_cols = st.multiselect("Group by columns", cols_avail)
            # rules
            if st.session_state["vb_agg_rules"]:
                st.markdown("Current rules")
                for i, r in enumerate(st.session_state["vb_agg_rules"]):
                    st.markdown(f"- {r['col']} → {', '.join(r['funcs'])}")
                    if st.button("Remove", key=f"agg_rm_{i}"):
                        st.session_state["vb_agg_rules"].pop(i)
                        st.success("Removed agg rule.")

            with st.form("agg_add_form", clear_on_submit=True):
                numeric_choices = [
                    c
                    for c in cols_avail
                    if pd.api.types.is_numeric_dtype(df_for_agg[c])
                ]
                sel = st.selectbox("Numeric column", ["--select--"] + numeric_choices)
                funcs = st.multiselect(
                    "Functions",
                    ["sum", "mean", "count", "min", "max", "median"],
                    default=["sum"],
                )
                add_agg = st.form_submit_button("Add rule")
                if add_agg:
                    if sel == "--select--" or not funcs:
                        st.error("Pick a numeric col and functions")
                    else:
                        st.session_state["vb_agg_rules"].append(
                            {"col": sel, "funcs": funcs}
                        )
                        st.success("Added agg rule")

            if st.button("Apply Aggregation"):
                try:
                    agg_map = {
                        r["col"]: r["funcs"] for r in st.session_state["vb_agg_rules"]
                    }
                    if grp_cols:
                        grouped = df_for_agg.groupby(grp_cols).agg(agg_map)
                        # flatten
                        grouped.columns = [
                            "_".join([str(c) for c in col]).strip()
                            if isinstance(col, tuple)
                            else str(col)
                            for col in grouped.columns.values
                        ]
                        agg_df = grouped.reset_index()
                    else:
                        tmp = df_for_agg.agg(agg_map)
                        agg_df = pd.DataFrame(tmp).T.reset_index()
                    st.session_state["_vb_agg_df"] = agg_df
                    st.success("Aggregation built; now available for charting")
                    st.dataframe(agg_df.head(50), use_container_width=True)
                except Exception as e:
                    st.error(f"Aggregation failed: {e}")

        # --------------------------
        # Smart Templates Tab
        # --------------------------
        with tabs[3]:
            st.markdown("## Analytics Templates")
            base_df = st.session_state.get(
                "_vb_agg_df",
                st.session_state.get(
                    "_vb_transformed_df",
                    st.session_state.get("_vb_refined_df", working_df.copy()),
                ),
            )

            st.markdown("### Time-series Summary (auto-detect date)")
            date_cols = [
                c
                for c in safe_cols(base_df)
                if pd.api.types.is_datetime64_any_dtype(base_df[c])
                or "date" in c.lower()
            ]
            if date_cols:
                dcol = st.selectbox("Date column", date_cols)
                resample = st.selectbox("Resample", ["D", "W", "M", "Q", "Y"], index=2)
                metric = st.selectbox(
                    "Metric (numeric)",
                    [
                        c
                        for c in safe_cols(base_df)
                        if pd.api.types.is_numeric_dtype(base_df[c])
                    ],
                    key="ts_metric",
                )
                if st.button("Generate Time-series Analysis"):
                    tmp = base_df.copy()
                    tmp[dcol] = pd.to_datetime(tmp[dcol], errors="coerce")
                    ts = (
                        tmp.set_index(dcol)
                        .resample(resample)[metric]
                        .sum()
                        .reset_index()
                    )
                    fig = px.line(
                        ts,
                        x=dcol,
                        y=metric,
                        title=f"{metric} over time ({resample})",
                        template="plotly_dark",
                    )
                    st.plotly_chart(fig, use_container_width=True)
            else:
                st.info(
                    "No date-like columns detected. Add or convert a column to datetime in Transform tab."
                )

            st.markdown("### Correlation Analysis")
            numeric_cols = [
                c
                for c in safe_cols(base_df)
                if pd.api.types.is_numeric_dtype(base_df[c])
            ]
            if len(numeric_cols) >= 2:
                if st.button("Show Correlation Matrix"):
                    corr = base_df[numeric_cols].corr()
                    fig = px.imshow(
                        corr,
                        text_auto=True,
                        title="Correlation Matrix",
                        template="plotly_dark",
                    )
                    st.plotly_chart(fig, use_container_width=True)
            else:
                st.info("At least two numeric columns required for correlation.")

            st.markdown("### Outlier Detection")
            val_col = st.selectbox(
                "Column for outlier detection", numeric_cols + [None], key="outlier_col"
            )
            method = st.selectbox(
                "Method", ["IQR", "Z-score"], index=0, key="outlier_method"
            )
            if st.button("Detect Outliers") and val_col:
                s = base_df[val_col].dropna()
                if method == "IQR":
                    q1 = s.quantile(0.25)
                    q3 = s.quantile(0.75)
                    iqr = q3 - q1
                    lower = q1 - 1.5 * iqr
                    upper = q3 + 1.5 * iqr
                    out = base_df[
                        (base_df[val_col] < lower) | (base_df[val_col] > upper)
                    ]
                else:
                    z = (s - s.mean()) / s.std()
                    mask = z.abs() > 3
                    out = base_df.loc[mask.index[mask], :]
                st.write(f"Found {len(out)} outliers")
                st.dataframe(out.head(50), use_container_width=True)

        # --------------------------
        # Charts Tab - advanced chart builder (lazy build)
        # --------------------------
        with tabs[4]:
            st.markdown("## Chart Builder (Advanced)")
            src_df = st.session_state.get(
                "_vb_agg_df",
                st.session_state.get(
                    "_vb_transformed_df",
                    st.session_state.get("_vb_refined_df", working_df.copy()),
                ),
            )
            st.caption(f"Chart source: {len(src_df)} rows, {len(src_df.columns)} cols")
            cols = safe_cols(src_df)

            chart_mode = st.selectbox(
                "Chart Mode", ["Single Chart", "Dual Axis", "Combo Chart"], index=0
            )
            chart_types = ["Bar", "Line", "Scatter", "Area", "Box", "Pie"]
            chart_kind = st.selectbox("Plot Type", chart_types, index=0)
            x_axis = st.selectbox("X Axis", ["--select--"] + cols)
            y_axis = st.selectbox(
                "Y Axis",
                ["--select--"]
                + [c for c in cols if pd.api.types.is_numeric_dtype(src_df[c])],
            )
            y2_axis = None
            if chart_mode in ("Dual Axis", "Combo Chart"):
                y2_axis = st.selectbox(
                    "Secondary Y Axis (numeric)",
                    ["--select--"]
                    + [c for c in cols if pd.api.types.is_numeric_dtype(src_df[c])],
                )
            color = st.selectbox(
                "Color / Group",
                ["None"]
                + [c for c in cols if not pd.api.types.is_numeric_dtype(src_df[c])],
                index=0,
            )
            add_trend = st.checkbox("Add trend line (linear)", value=False)
            title = st.text_input("Title", value=f"{chart_kind} Chart")
            subtitle = st.text_input("Subtitle (for report)", value="")
            description = st.text_area("Description (for PDF/PPT)", value="")
            # lazy build: user presses Build Chart
            if st.button("Build Chart"):
                try:
                    fig = None
                    kwargs = {}
                    if color and color != "None":
                        kwargs["color"] = color
                    if chart_kind == "Bar":
                        fig = px.bar(
                            src_df,
                            x=x_axis if x_axis != "--select--" else None,
                            y=y_axis if y_axis != "--select--" else None,
                            **kwargs,
                            title=title,
                            template="plotly_dark",
                        )
                    elif chart_kind == "Line":
                        fig = px.line(
                            src_df,
                            x=x_axis if x_axis != "--select--" else None,
                            y=y_axis if y_axis != "--select--" else None,
                            **kwargs,
                            title=title,
                            template="plotly_dark",
                        )
                    elif chart_kind == "Scatter":
                        fig = px.scatter(
                            src_df,
                            x=x_axis if x_axis != "--select--" else None,
                            y=y_axis if y_axis != "--select--" else None,
                            **kwargs,
                            title=title,
                            template="plotly_dark",
                        )
                    elif chart_kind == "Area":
                        fig = px.area(
                            src_df,
                            x=x_axis if x_axis != "--select--" else None,
                            y=y_axis if y_axis != "--select--" else None,
                            **kwargs,
                            title=title,
                            template="plotly_dark",
                        )
                    elif chart_kind == "Box":
                        fig = px.box(
                            src_df,
                            x=x_axis if x_axis != "--select--" else None,
                            y=y_axis if y_axis != "--select--" else None,
                            title=title,
                            template="plotly_dark",
                        )
                    elif chart_kind == "Pie":
                        fig = px.pie(
                            src_df,
                            names=x_axis if x_axis != "--select--" else None,
                            values=y_axis if y_axis != "--select--" else None,
                            title=title,
                            template="plotly_dark",
                        )

                    # add secondary axis if requested (compose as go.Figure)
                    if (
                        chart_mode in ("Dual Axis", "Combo Chart")
                        and y2_axis
                        and y2_axis != "--select--"
                    ):
                        # build a combined figure
                        go_fig = go.Figure()
                        # primary trace
                        if fig:
                            for trace in fig.data:
                                go_fig.add_trace(trace)
                        # add y2 trace
                        if chart_kind in ("Bar", "Combo Chart"):
                            go_fig.add_trace(
                                go.Scatter(
                                    x=src_df[x_axis],
                                    y=src_df[y2_axis],
                                    name=y2_axis,
                                    yaxis="y2",
                                )
                            )
                        else:
                            go_fig.add_trace(
                                go.Scatter(
                                    x=src_df[x_axis],
                                    y=src_df[y2_axis],
                                    name=y2_axis,
                                    yaxis="y2",
                                )
                            )
                        go_fig.update_layout(
                            yaxis=dict(title=y_axis),
                            yaxis2=dict(title=y2_axis, overlaying="y", side="right"),
                            title=title,
                            template="plotly_dark",
                        )
                        fig = go_fig

                    # add trendline overlay if requested
                    if add_trend and chart_kind in ("Scatter", "Line", "Bar"):
                        try:
                            # simple OLS
                            tmp = src_df[[x_axis, y_axis]].dropna()
                            if pd.api.types.is_numeric_dtype(tmp[x_axis]):
                                coeffs = np.polyfit(tmp[x_axis], tmp[y_axis], 1)
                                trend_y = np.polyval(coeffs, tmp[x_axis])
                                fig.add_trace(
                                    go.Scatter(
                                        x=tmp[x_axis],
                                        y=trend_y,
                                        mode="lines",
                                        name="Trend",
                                        line=dict(dash="dash"),
                                    )
                                )
                            else:
                                # attempt numeric index trend
                                tmp_idx = np.arange(len(tmp))
                                coeffs = np.polyfit(tmp_idx, tmp[y_axis].values, 1)
                                trend_y = np.polyval(coeffs, tmp_idx)
                                fig.add_trace(
                                    go.Scatter(
                                        x=tmp[x_axis],
                                        y=trend_y,
                                        mode="lines",
                                        name="Trend",
                                        line=dict(dash="dash"),
                                    )
                                )
                        except Exception:
                            st.warning(
                                "Trendline generation failed; ensure numeric axes."
                            )

                    # finalize & render
                    fig.update_layout(paper_bgcolor="rgba(0,0,0,0)", height=520)
                    st.plotly_chart(fig, use_container_width=True)

                    # actions: add to composer, single pdf, add text block, add table
                    c1, c2, c3, c4 = st.columns([1, 1, 1, 1])
                    with c1:
                        if st.button("➕ Add Chart to Composer"):
                            # export image
                            img_bytes = None
                            try:
                                img_bytes = fig.to_image(format="png")
                            except Exception:
                                try:
                                    import plotly.io as pio

                                    img_bytes = pio.to_image(fig, format="png")
                                except Exception as e:
                                    st.error(f"Export failed: {e}")
                                    img_bytes = None
                            img_path = write_temp_png_from_bytes(img_bytes)
                            st.session_state["my_composer"].append(
                                {
                                    "id": uuid.uuid4().hex,
                                    "type": "chart",
                                    "title": title,
                                    "subtitle": subtitle,
                                    "description": description,
                                    "image": img_bytes,
                                    "image_path": img_path,
                                    "created_at": datetime.now().isoformat(),
                                }
                            )
                            st.success("Added to Composer")
                    with c2:
                        if st.button("Download Single Chart -> PDF"):
                            # single chart PDF
                            if "fig" in locals() and fig:
                                try:
                                    img_bytes = fig.to_image(format="png")
                                except Exception:
                                    import plotly.io as pio

                                    img_bytes = pio.to_image(fig, format="png")
                                b64 = base64.b64encode(img_bytes).decode("utf-8")
                                html = f"<html><body style='font-family:Arial;padding:20px;'><h1>{title}</h1><p>{description}</p><img src='data:image/png;base64,{b64}' style='max-width:100%;'/></body></html>"
                                ok, pdf_bytes = PDFEngine.generate(
                                    html, f"chart_{uuid.uuid4().hex}.pdf"
                                )
                                if ok:
                                    st.download_button(
                                        "Download PDF",
                                        pdf_bytes,
                                        file_name=f"{title}.pdf",
                                        mime="application/pdf",
                                    )
                                else:
                                    st.error("PDF generation failed.")
                    with c3:
                        if st.button("Add Text Block to Composer"):
                            st.session_state["my_composer"].append(
                                {
                                    "id": uuid.uuid4().hex,
                                    "type": "text",
                                    "title": title,
                                    "subtitle": subtitle,
                                    "description": description,
                                    "created_at": datetime.now().isoformat(),
                                }
                            )
                            st.success("Added text block")
                    with c4:
                        if st.button("Add Data Table to Composer"):
                            st.session_state["my_composer"].append(
                                {
                                    "id": uuid.uuid4().hex,
                                    "type": "table",
                                    "title": f"Table: {title}",
                                    "table_df": src_df.head(200).copy(),
                                    "description": description,
                                    "created_at": datetime.now().isoformat(),
                                }
                            )
                            st.success("Added table (top 200 rows) to composer")

                except Exception as e:
                    st.error(f"Chart build failed: {e}")

        # --------------------------
        # Composer Tab - reorder, edit, export
        # --------------------------
        with tabs[5]:
            st.markdown("## Composer (Report Designer)")
            items = st.session_state["my_composer"]
            if not items:
                st.info(
                    "Composer empty — build charts / tables / text blocks from 'Charts' tab."
                )
            else:
                # show list with move up/down/remove / edit
                for i, it in enumerate(list(items)):
                    cols = st.columns([1, 4, 1, 1])
                    with cols[0]:
                        if it.get("type") == "chart" and it.get("image"):
                            try:
                                b64 = base64.b64encode(it["image"]).decode("utf-8")
                                st.image(f"data:image/png;base64,{b64}", width=140)
                            except Exception:
                                st.text("preview")
                        elif it.get("type") == "table":
                            st.text("Table")
                        else:
                            st.text("Text")
                    with cols[1]:
                        st.markdown(f"**{it.get('title', 'Untitled')}**")
                        if it.get("subtitle"):
                            st.markdown(f"_{it.get('subtitle')}_")
                        if it.get("description"):
                            st.markdown(it.get("description"))
                    with cols[2]:
                        if st.button("Up", key=f"comp_up_{i}") and i > 0:
                            (
                                st.session_state["my_composer"][i],
                                st.session_state["my_composer"][i - 1],
                            ) = (
                                st.session_state["my_composer"][i - 1],
                                st.session_state["my_composer"][i],
                            )
                            st.success("Moved up")
                        if (
                            st.button("Down", key=f"comp_down_{i}")
                            and i < len(items) - 1
                        ):
                            (
                                st.session_state["my_composer"][i],
                                st.session_state["my_composer"][i + 1],
                            ) = (
                                st.session_state["my_composer"][i + 1],
                                st.session_state["my_composer"][i],
                            )
                            st.success("Moved down")
                    with cols[3]:
                        if st.button("Remove", key=f"comp_rm_{i}"):
                            st.session_state["my_composer"].pop(i)
                            st.success("Removed from composer")

                st.markdown("---")
                # Export controls
                exp1, exp2, exp3 = st.columns([1, 1, 1])
                with exp1:
                    if st.button("Export Composer -> PDF"):
                        ok, pdf_or_err = build_pdf_from_composer(
                            items, title="Composer Report", theme="professional"
                        )
                        if not ok:
                            st.error(f"PDF failed: {pdf_or_err}")
                        else:
                            st.download_button(
                                "Download PDF",
                                pdf_or_err,
                                file_name="composer_report.pdf",
                                mime="application/pdf",
                            )
                with exp2:
                    if st.button("Export Composer -> PPTX"):
                        ok, out = build_pptx_from_composer(
                            items, title="Composer Report"
                        )
                        if not ok:
                            st.error(f"PPTX failed: {out}")
                        else:
                            st.download_button(
                                "Download PPTX",
                                out,
                                file_name="composer_report.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            )
                with exp3:
                    if st.button("Email Composer (attach PDF)"):
                        # create pdf then email using smtp_config
                        ok, pdf_or_err = build_pdf_from_composer(
                            items, title="Composer Report"
                        )
                        if not ok:
                            st.error(f"PDF build failed: {pdf_or_err}")
                        else:
                            smtp_cfg = st.session_state.get("smtp_config", {})
                            if not smtp_cfg.get("host"):
                                st.info(
                                    "No SMTP configured. PDF ready to download instead."
                                )
                                st.download_button(
                                    "Download PDF",
                                    pdf_or_err,
                                    file_name="composer_report.pdf",
                                    mime="application/pdf",
                                )
                            else:
                                recip = st.text_input(
                                    "Recipient email",
                                    value=smtp_cfg.get(
                                        "default_recipient", "manager@company.com"
                                    ),
                                    key="email_recipient",
                                )
                                if st.button("Send Email Now"):
                                    try:
                                        from email.message import EmailMessage
                                        import smtplib

                                        msg = EmailMessage()
                                        msg["Subject"] = "Composer Report"
                                        msg["From"] = smtp_cfg.get(
                                            "from",
                                            smtp_cfg.get(
                                                "user", "no-reply@example.com"
                                            ),
                                        )
                                        msg["To"] = recip
                                        msg.set_content(
                                            "Please find attached the Composer report."
                                        )
                                        msg.add_attachment(
                                            pdf_or_err,
                                            maintype="application",
                                            subtype="pdf",
                                            filename="composer_report.pdf",
                                        )
                                        s = smtplib.SMTP(
                                            smtp_cfg["host"], smtp_cfg.get("port", 587)
                                        )
                                        s.starttls()
                                        if smtp_cfg.get("user"):
                                            s.login(
                                                smtp_cfg.get("user"),
                                                smtp_cfg.get("pass"),
                                            )
                                        s.send_message(msg)
                                        s.quit()
                                        st.success("Email sent.")
                                    except Exception as e:
                                        st.error(f"Send failed: {e}")

    # final footer
    ThemeEngine.render_comment_section("visual_builder_pro")


# --- SEARCH ENGINE PAGE ---
def page_search():
    ThemeEngine.render_header("Global Registry", "Corporate Intelligence Search")

    df = st.session_state.company_df

    # Search Interface
    st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
    c1, c2 = st.columns([3, 1])
    with c1:
        query = st.text_input(
            "Search Corp ID, Name, or Tax Number",
            placeholder="e.g., COMP-1002, Wayne Enterprises, Stark",
            help="You can search for multiple entities separated by commas.",
        )
    with c2:
        st.markdown("<div style='margin-top:28px'></div>", unsafe_allow_html=True)
        search_btn = st.button(
            "🔍 Search Database", type="primary", use_container_width=True
        )

    st.markdown("</div>", unsafe_allow_html=True)

    if search_btn or query:
        # Search Logic (Multi-Entity Support)
        if query:
            terms = [t.strip() for t in query.split(",") if t.strip()]
            mask = pd.Series([False] * len(df))
            for term in terms:
                term_mask = df["company_id"].astype(str).str.contains(
                    term, case=False
                ) | df["company_name"].astype(str).str.contains(term, case=False)
                mask = mask | term_mask
            results = df[mask]
        else:
            results = pd.DataFrame()  # Empty if no query

        if results.empty and query:
            st.warning("No records found matching your criteria.")

        elif not results.empty:
            # --- COMPARISON / SELECTION VIEW ---
            if len(results) > 1:
                st.markdown(f"#### 🔍 Found {len(results)} matches")
                st.info("Select an entity from the list below to view full details.")

                # Comparison Table
                st.dataframe(
                    results[
                        [
                            "company_id",
                            "company_name",
                            "sector",
                            "status",
                            "revenue_mm",
                            "risk_score",
                        ]
                    ],
                    use_container_width=True,
                    hide_index=True,
                )

                # Selection for Details
                selected_id_compare = st.selectbox(
                    "View Details For:",
                    results["company_id"].tolist(),
                    format_func=lambda x: f"{x} - {df[df['company_id'] == x]['company_name'].iloc[0]}",
                    key="multi_search_select",
                )
                entity = df[df["company_id"] == selected_id_compare].iloc[0]
                show_details = True
            else:
                # Direct Match
                entity = results.iloc[0]
                show_details = True

            if show_details:
                # --- DETAILED VIEW ---
                st.markdown("---")

                # Header
                st.markdown(
                    f"""
                <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:20px">
                    <div>
                        <h1 style="margin:0">{entity["company_name"]}</h1>
                        <div style="color:var(--primary-color); font-family:'JetBrains Mono'">{entity["company_id"]}</div>
                    </div>
                    <div style="text-align:right">
                        <div style="background:{"#00e396" if entity["status"] == "Active" else "#ff0055"}; padding:5px 15px; border-radius:20px; color:#fff; font-weight:bold; display:inline-block">
                            {entity["status"].upper()}
                        </div>
                    </div>
                </div>
                """,
                    unsafe_allow_html=True,
                )

                # KPIs
                k1, k2, k3, k4 = st.columns(4)
                with k1:
                    st.markdown(
                        ThemeEngine.render_kpi_card(
                            "Annual Revenue", f"€{entity['revenue_mm']:,.1f}M"
                        ),
                        unsafe_allow_html=True,
                    )
                with k2:
                    st.markdown(
                        ThemeEngine.render_kpi_card(
                            "Global Headcount", f"{entity['employees']:,}"
                        ),
                        unsafe_allow_html=True,
                    )
                with k3:
                    st.markdown(
                        ThemeEngine.render_kpi_card(
                            "Founded", str(entity["founded_year"])
                        ),
                        unsafe_allow_html=True,
                    )
                with k4:
                    st.markdown(
                        ThemeEngine.render_kpi_card(
                            "Risk Score",
                            f"{entity['risk_score']}/100",
                            delta="-2" if entity["risk_score"] > 50 else "+1",
                            delta_desc="vs Sector Avg",
                        ),
                        unsafe_allow_html=True,
                    )

                # Modules
                tab1, tab2, tab3 = st.tabs(
                    [
                        "📜 Corporate History",
                        "🕸 Entity Hierarchy",
                        "📊 Financial Health",
                    ]
                )

                with tab1:
                    st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
                    st.markdown("#### Timeline of Events")
                    events = [
                        {
                            "date": "2024-01-15",
                            "type": "Filing",
                            "desc": "Annual Report Filed (10-K)",
                        },
                        {
                            "date": "2023-11-20",
                            "type": "M&A",
                            "desc": "Acquired subsidiary technology unit",
                        },
                        {
                            "date": "2023-05-10",
                            "type": "Legal",
                            "desc": "Settlement reached in IP dispute",
                        },
                        {
                            "date": "2022-08-01",
                            "type": "Executive",
                            "desc": "New CFO appointed",
                        },
                    ]
                    for event in events:
                        st.markdown(
                            f"""
                        <div style="border-left:2px solid var(--primary-color); padding-left:15px; margin-bottom:15px">
                            <div style="color:var(--muted-color); font-size:0.8rem">{event["date"]} | {event["type"]}</div>
                            <div style="font-weight:600">{event["desc"]}</div>
                        </div>
                        """,
                            unsafe_allow_html=True,
                        )
                    st.markdown("</div>", unsafe_allow_html=True)

                with tab2:
                    c_hier, c_desc = st.columns([2, 1])
                    with c_hier:
                        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
                        st.markdown("#### Control Network")

                        # Mock Hierarchy for Visualization
                        # Create a small tree around the entity

                        fam_data = {
                            "ids": [
                                entity["company_name"],
                                "North America Div",
                                "Europe Div",
                                "Asia Ops",
                                "R&D Lab",
                                "Sales Corp",
                            ],
                            "parents": [
                                "",
                                entity["company_name"],
                                entity["company_name"],
                                entity["company_name"],
                                "North America Div",
                                "Europe Div",
                            ],
                            "values": [
                                entity["revenue_mm"],
                                entity["revenue_mm"] * 0.4,
                                entity["revenue_mm"] * 0.3,
                                entity["revenue_mm"] * 0.3,
                                entity["revenue_mm"] * 0.1,
                                entity["revenue_mm"] * 0.1,
                            ],
                        }

                        fig_tree = go.Figure(
                            go.Treemap(
                                labels=fam_data["ids"],
                                parents=fam_data["parents"],
                                values=fam_data["values"],
                                root_color="rgba(0,0,0,0)",
                            )
                        )
                        fig_tree.update_layout(
                            template="plotly_dark",
                            margin=dict(t=0, l=0, r=0, b=0),
                            height=400,
                        )
                        st.plotly_chart(fig_tree, use_container_width=True)
                        st.markdown("</div>", unsafe_allow_html=True)

                    with c_desc:
                        st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
                        st.info(
                            f"Ultimate Beneficial Owner identified as Holding Group Alpha (75%)."
                        )
                        st.write(f"Subsidiaries: {entity['subsidiaries_count']}")
                        st.write(
                            f"Parent ID: {entity['parent_id'] if entity['parent_id'] else 'None (Top Co)'}"
                        )
                        st.button("Download Entity Report", use_container_width=True)
                        st.markdown("</div>", unsafe_allow_html=True)

                with tab3:
                    st.markdown("<div class='glass-card'>", unsafe_allow_html=True)
                    st.line_chart(np.random.randn(20, 3) + 10)
                    st.markdown("</div>", unsafe_allow_html=True)

    # Comment Section Integration
    ThemeEngine.render_comment_section(
        f"search_{entity['company_id'] if 'entity' in locals() else 'general'}"
    )


# --- SCHEDULER PAGE ---
def page_scheduler():
    ThemeEngine.render_header("Automation Hub", "Task Scheduler")

    st.markdown("<div class='glass-card'>", unsafe_allow_html=True)

    c1, c2 = st.columns([1, 2])

    with c1:
        st.markdown("#### New Schedule")

        # Check for pre-filled data
        defaults = st.session_state.get("scheduler_defaults", {})
        def_name = defaults.get("name", "Weekly Portfolio Summary")

        job_name = st.text_input("Job Name", def_name)

        # Source Selection
        source_options = ["Standard Report", "Saved View"]
        def_source_idx = 0
        if defaults.get("source") == "Saved View":
            def_source_idx = 1

        source_type = st.radio("Source Type", source_options, index=def_source_idx)

        selected_view = None
        if source_type == "Saved View":
            if st.session_state.saved_views:
                # If passed from Visual Builder, try to pre-select
                pre_view_name = defaults.get("view_config", None)
                view_opts = list(st.session_state.saved_views.keys())

                # If the passed name matches a key, use it; else default 0
                view_idx = 0
                if pre_view_name in view_opts:
                    view_idx = view_opts.index(pre_view_name)

                selected_view = st.selectbox(
                    "Select Saved View", view_opts, index=view_idx
                )
            else:
                st.warning(
                    "No saved views available. Please create one in Visual Builder."
                )

        freq = st.selectbox("Frequency", ["Daily", "Weekly", "Monthly", "Quarterly"])
        time_val = st.time_input("Run Time", value=datetime.now().time())
        recipients = st.text_area("Recipients", "admin@aurora.io, management@aurora.io")

        fmt_idx = 0
        if defaults.get("format") == "PDF Report":
            fmt_idx = 0

        format_type = st.radio(
            "Format", ["PDF Report", "Excel Dump", "JSON API Payload"], index=fmt_idx
        )

        if st.button("Create Schedule", type="primary"):
            job_details = {
                "id": uuid.uuid4().hex[:8],
                "name": job_name,
                "freq": f"{freq} @ {time_val.strftime('%H:%M')}",
                "recipients": recipients.split(","),
                "format": format_type,
                "status": "Active",
                "next_run": (datetime.now() + timedelta(days=7)).strftime("%Y-%m-%d"),
                "source": source_type,
            }
            if selected_view:
                job_details["view_linked"] = selected_view

            if "scheduled_jobs" not in st.session_state:
                st.session_state.scheduled_jobs = []
            st.session_state.scheduled_jobs.append(job_details)

            st.success("Job Scheduled!")
            # Clear defaults after usage
            if "scheduler_defaults" in st.session_state:
                del st.session_state["scheduler_defaults"]

    with c2:
        st.markdown("#### Active Schedules")
        if "scheduled_jobs" in st.session_state and st.session_state.scheduled_jobs:
            jobs_df = pd.DataFrame(st.session_state.scheduled_jobs)
            st.dataframe(
                jobs_df[["name", "freq", "format", "status", "next_run"]],
                use_container_width=True,
                hide_index=True,
            )
        else:
            st.info("No active schedules found.")

    st.markdown("</div>", unsafe_allow_html=True)


# ==============================================================================
# 8. ROUTING & ENTRY POINT
# ==============================================================================


def main_router():
    # Sidebar Navigation
    with st.sidebar:
        st.markdown(
            """
        <div style="text-align:center; margin-bottom:20px;">
            <div style="font-size:3rem;">'💠'</div>
            <h2 style="margin:0;">AURORA</h2>
            <div style="color:var(--muted-color); font-size:0.8rem;">ENTERPRISE v8.0</div>
        </div>
        """,
            unsafe_allow_html=True,
        )

        menu = {
            "🏠 Dashboard": page_dashboard,
            "📊 Analytics": page_analytics,
            "💹 Trading": page_trading,
            "📈 Revenue": page_sales,
            "💾 Data Studio": page_data_wrangling,
            "🎨 Visual Builder": page_visual_builder,
            "🔍 Search": page_search,
            "📄 Reports": page_reports,
            "📅 Scheduler": page_scheduler,
            "⚙️ Settings": page_settings,
        }

        selected = st.radio("MODULES", list(menu.keys()), label_visibility="collapsed")

        st.markdown("---")

        # Sidebar Widgets
        st.markdown(
            "<div style='font-size:0.85rem; font-weight:600; color:var(--muted-color); margin-bottom:10px'>MARKET STATUS</div>",
            unsafe_allow_html=True,
        )

        # Mini Ticker
        mkt = st.session_state.market_df.iloc[-1]
        col_tick = "#00e396" if mkt["returns_pct"] > 0 else "#ff0055"

        # Pre-calc values to avoid f-string syntax errors
        price_val = mkt["price"]
        ret_val = mkt["returns_pct"] * 100

        # Ticker HTML Construction
        p_str = "{:.2f}".format(price_val)
        r_str = "{:.2f}%".format(ret_val)

        html_ticker = (
            "<div style='background:rgba(255,255,255,0.05); padding:10px; border-radius:8px;'>"
            "<div style='display:flex; justify-content:space-between;'>"
            "<span>S&P 500 (Syn)</span>"
            "<span style='color:" + col_tick + "'>" + p_str + "</span>"
            "</div>"
            "<div style='font-size:0.8rem; text-align:right; color:"
            + col_tick
            + "'>"
            + r_str
            + "</div>"
            "</div>"
        )

        st.markdown(html_ticker, unsafe_allow_html=True)

        st.markdown("<div style='margin-top:20px'></div>", unsafe_allow_html=True)
        if st.button("Logout", use_container_width=True):
            st.session_state.user = None
            st.rerun()

    # Routing
    if selected in menu:  #
        try:
            menu[selected]()
        except Exception as e:
            st.error(f"Critical Module Error: {str(e)}")
            st.exception(e)


if __name__ == "__main__":
    main_router()
