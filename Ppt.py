from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def add_stylish_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 55, 95)
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1))
    sf = subtitle_box.text_frame
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(210, 230, 255)

def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(232, 242, 250)
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 80, 140)

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 250, 255)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 90, 150)

    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8.5), Inches(4))
    cf = content_box.text_frame
    for bullet in bullets:
        bp = cf.add_paragraph()
        bp.text = bullet
        bp.font.size = Pt(22)
        bp.font.color.rgb = RGBColor(40, 40, 40)
        bp.level = 0

# Build deck
add_stylish_title_slide("Migration Support Provided", "Cross-Team Platform Migration Overview")

add_section_slide("Executive Summary")

add_content_slide("Executive Summary", [
    "Our team provided complete migration and technical support to the partner team.",
    "Ensured smooth transition with validation, troubleshooting, and platform alignment.",
    "Post‑migration data staleness observed due to missing code deployment.",
    "Action plan identified to resolve data refresh pipeline issues."
])

add_section_slide("Migration Scope")

add_content_slide("Migration Scope", [
    "Migrated dashboards, data sources, scripts, and platform components.",
    "Aligned platform settings and environment configurations.",
    "Performed dependency validation and end‑to‑end testing.",
    "Completed access provisioning and platform onboarding."
])

add_section_slide("Support Provided")

add_content_slide("Support Provided by Our Team", [
    "Code review, refactoring, and onboarding support.",
    "Platform troubleshooting & configuration assistance.",
    "Daily sync meetings, KT sessions, walkthroughs.",
    "Documentation handover and smooth coordination."
])

add_section_slide("Post‑Migration Issue")

add_content_slide("Issue Found After Migration", [
    "Table status indicates stale data — new data is not loading.",
    "Root Cause: No new code push from partner team post‑migration.",
    "Schedulers/ETL not reconnected in the new platform.",
    "Impact: Reports showing outdated data → affects decision‑making."
])

add_section_slide("Action Plan")

add_content_slide("Action Plan / Recommendations", [
    "Partner team to push pending code changes immediately.",
    "Reconnect ETL/schedulers and validate end‑to‑end run status.",
    "Enable monitoring alerts for pipeline health checks.",
    "Complete validation cycle & close migration formally."
])

file_path = "/mnt/data/NEXT_GEN_MIGRATION_PRESENTATION.pptx"
prs.save(file_path)

file_path
