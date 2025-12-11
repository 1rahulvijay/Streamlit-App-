"""
Generate an enterprise-style PPT that matches the provided sample layout.

Produces: Exact_Style_Migration_Presentation.pptx

Prereqs:
    pip install python-pptx matplotlib pillow
Place assets (icons/screenshots) in the same folder or update the paths:
    - icon_check.png
    - icon_tool.png
    - icon_guidance.png
    - icon_question.png
    - logo.png
    - before.png
    - after.png
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
import matplotlib.pyplot as plt
import os

# --- Config / content (edit these variables) ---
TEAM_NAME = "Platform Engineering Team"
OTHER_TEAM = "Data Analytics Team"
PRESENTER = "Your Name"
DATE = "2025-12-11"

# Metrics (fill your actual numbers)
METRICS = {
    "Scripts migrated": 18,
    "Dashboards moved": 12,
    "Data sources configured": 9,
    "Hours invested": 120,
    "Issues resolved": 34
}

# Asset filenames (replace with your icons/screenshots)
ICON_CHECK = "icon_check.png"
ICON_TOOL = "icon_tool.png"
ICON_GUIDE = "icon_guidance.png"
ICON_Q = "icon_question.png"
LOGO = "logo.png"
BEFORE_IMG = "before.png"
AFTER_IMG = "after.png"

OUT_FILE = "Exact_Style_Migration_Presentation.pptx"


# --- Helpers ---
def set_cell_text(cell, text, bold=False, size=12, color=RGBColor(255,255,255), align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

def add_rounded_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, radius=0.1):
    # Use simple rectangle (python-pptx doesn't have explicit rounded rect by default in older versions)
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb
    else:
        shp.line.fill.background()
    return shp

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(12, 47, 84)  # deep navy

    # Left gear icon area (big visual)
    left_panel = add_rounded_rect(slide, Inches(0.4), Inches(0.6), Inches(4.2), Inches(4.2), RGBColor(22, 78, 130))
    left_panel.fill.fore_color.rgb = RGBColor(20, 70, 120)

    # Title text
    title_box = slide.shapes.add_textbox(Inches(4.8), Inches(0.9), Inches(5.2), Inches(1.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MIGRATION SUPPORT PROVIDER"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    sub_box = slide.shapes.add_textbox(Inches(4.8), Inches(1.9), Inches(5.2), Inches(1.0))
    sf = sub_box.text_frame
    p2 = sf.paragraphs[0]
    p2.text = f"Provided by {TEAM_NAME}\nCollaboration with {OTHER_TEAM} – Migration to [PLATFORM/SOLUTION]"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(210, 230, 245)

    # presenter & date bottom left
    pd_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(4.8), Inches(0.6))
    pdf = pd_box.text_frame
    pdf.text = f"{PRESENTER} | {DATE}"
    pdf.paragraphs[0].font.color.rgb = RGBColor(180, 200, 220)

    # Logo bottom right (if exists)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(8.5), Inches(5.0), width=Inches(1.4))


def add_executive_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # pale blue background panel
    panel = add_rounded_rect(slide, Inches(0.4), Inches(0.4), prs.slide_width - Inches(0.8), Inches(2.0), RGBColor(37, 121, 163))
    panel.fill.fore_color.rgb = RGBColor(44, 134, 180)
    # White banner top-left "EXECUTIVE SUMMARY"
    banner = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(3.2), Inches(0.5))
    bn = banner.text_frame
    bn.text = "EXECUTIVE SUMMARY"
    bn.paragraphs[0].font.bold = True
    bn.paragraphs[0].font.size = Pt(14)
    bn.paragraphs[0].font.color.rgb = RGBColor(255,255,255)

    # Checklist bullets w/ mini-icons
    x = Inches(0.8)
    y = Inches(1.1)
    bullet_texts = [
        "Your team provided end-to-end support to all migration items on-time.",
        "Post-migration monitoring revealed data staleness issues due to missing code push.",
        "Action plan defined to resolve data issues and stabilize operations."
    ]
    for i, txt in enumerate(bullet_texts):
        # small circle icon - use shapes to simulate
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(i*0.45), Inches(0.28), Inches(0.28))
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor(34, 197, 106)  # green
        circ.line.fill.background()
        # text
        tb = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(i*0.45), Inches(8), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(245,245,245)

    # mini flow icons row (illustrative)
    # Add three small circles connected with lines to mimic the sample
    fx = Inches(6.4)
    fy = Inches(1.1)
    for i in range(3):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, fx + Inches(i*0.9), fy, Inches(0.5), Inches(0.5))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(255,255,255)
        c.line.fill.background()


def add_migration_scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # background soft blue
    bg = add_rounded_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RGBColor(13, 52, 88))
    bg.fill.fore_color.rgb = RGBColor(10, 45, 80)

    # title bar
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(5.5), Inches(0.9))
    t = title_box.text_frame
    t.text = "MIGRATION SCOPE"
    t.paragraphs[0].font.size = Pt(22)
    t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.color.rgb = RGBColor(255,255,255)

    # left bullet column inside lighter panel
    panel = add_rounded_rect(slide, Inches(0.6), Inches(1.4), Inches(9.0), Inches(2.3), RGBColor(23, 106, 150))
    panel.fill.fore_color.rgb = RGBColor(28, 116, 165)

    bullets = [
        f"{OTHER_TEAM} required migration to [Platform/Solution].",
        "Goal: Standardize platform usage, improve maintainability, reduce tech debt.",
        "Your team acted as technical owners and migration partners.",
        "Performed dependency validation, access mapping and KT sessions."
    ]

    for idx, b in enumerate(bullets):
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.6 + idx*0.5), Inches(8.4), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(235,245,255)


def add_support_provided_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # darker panel
    main = add_rounded_rect(slide, Inches(0.4), Inches(0.4), prs.slide_width - Inches(0.8), Inches(2.6), RGBColor(12, 35, 58))
    main.fill.fore_color.rgb = RGBColor(12,35,58)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(5), Inches(0.6))
    title_box.text_frame.text = "SUPPORT PROVIDED BY OUR TEAM"
    title_box.text_frame.paragraphs[0].font.size = Pt(18)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(240,240,240)

    # three columns: Technical Support | Collaboration & Guidance | What Worked Well
    col_x = [Inches(0.6), Inches(3.6), Inches(6.6)]
    headings = ["TECHNICAL SUPPORT", "COLLABORATION & GUIDANCE", "WHAT WORKED WELL"]
    columns = [
        ["Code review and refactoring", "ETL/API validation & reconnections", "Platform configuration"],
        ["Daily sync meetings", "Live walkthroughs & KT sessions", "Documentation handover"],
        ["No major outages during cutover", "Dependencies validated", "Onboarding completed"]
    ]
    for i in range(3):
        # heading box
        hb = slide.shapes.add_textbox(col_x[i], Inches(0.95), Inches(2.6), Inches(0.4))
        hb_tf = hb.text_frame
        hb_tf.text = headings[i]
        hb_tf.paragraphs[0].font.size = Pt(12)
        hb_tf.paragraphs[0].font.bold = True
        hb_tf.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)
        # items
        for j, it in enumerate(columns[i]):
            tb = slide.shapes.add_textbox(col_x[i], Inches(1.4 + j*0.45), Inches(2.6), Inches(0.4))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + it
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(220,235,245)


def add_statistics_slide(prs, metrics):
    # create a matplotlib chart and embed
    labels = list(metrics.keys())
    values = list(metrics.values())

    fig, ax = plt.subplots(figsize=(6,3))
    bars = ax.bar(labels, values)
    ax.set_ylim(0, max(values)*1.4)
    ax.set_ylabel("Count / Hours")
    ax.set_title("Statistics / Metrics")
    plt.xticks(rotation=20, ha='right')
    plt.tight_layout()
    chart_path = "metrics_chart.png"
    fig.savefig(chart_path, dpi=150)
    plt.close(fig)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # dark navy background
    panel = add_rounded_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RGBColor(8, 27, 45))
    panel.fill.fore_color.rgb = RGBColor(8,27,45)

    # left: chart
    slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.0), width=Inches(5.6), height=Inches(3.0))

    # right: root cause text block
    rc_box = add_rounded_rect(slide, Inches(6.6), Inches(1.0), Inches(3.1), Inches(3.0), RGBColor(24, 64, 91))
    rc_box.fill.fore_color.rgb = RGBColor(24,64,91)
    tb = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(2.8), Inches(2.8))
    tf = tb.text_frame
    tf.text = "ROOT CAUSE\n"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(255,220,180)
    p = tf.add_paragraph()
    p.text = "- No new code push from consuming team\n- ETL/schedulers not reconnected\n- Final integration steps not completed"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(220,240,255)


def add_issue_found_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # split background two-tone
    left = add_rounded_rect(slide, Inches(0.4), Inches(0.6), Inches(5.8), Inches(3.6), RGBColor(195, 48, 48))
    left.fill.fore_color.rgb = RGBColor(195,48,48)  # red
    right = add_rounded_rect(slide, Inches(6.3), Inches(0.6), Inches(3.2), Inches(3.6), RGBColor(244, 165, 66))
    right.fill.fore_color.rgb = RGBColor(244,165,66)  # orange

    # left: problem header + root cause
    hbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(5.4), Inches(0.6))
    hbox.text_frame.text = "PROBLEM: Table status shows data is stale — new data not getting loaded."
    hbox.text_frame.paragraphs[0].font.size = Pt(12)
    hbox.text_frame.paragraphs[0].font.bold = True
    hbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)

    rc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.2), Inches(2.4))
    rc_tf = rc_box.text_frame
    rc_tf.text = "ROOT CAUSE"
    rc_tf.paragraphs[0].font.bold = True
    rc_tf.paragraphs[0].font.size = Pt(12)
    rc_tf.paragraphs[0].font.color.rgb = RGBColor(255,235,235)
    p = rc_tf.add_paragraph()
    p.text = "- No new code push from the consuming team\n- ETL jobs were not re-attached after cutover\n- Ownership handoff incomplete"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255,235,235)

    # right: impact list
    im_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.4), Inches(2.8), Inches(2.4))
    im_tf = im_box.text_frame
    im_tf.text = "IMPACT"
    im_tf.paragraphs[0].font.bold = True
    im_tf.paragraphs[0].font.size = Pt(12)
    im_tf.paragraphs[0].font.color.rgb = RGBColor(40,20,0)
    p2 = im_tf.add_paragraph()
    p2.text = "- Data accuracy affected\n- Reports show outdated values\n- Business decisions impacted"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(40,20,0)


def add_action_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Before & After cards on left
    before_card = add_rounded_rect(slide, Inches(0.6), Inches(0.6), Inches(4.5), Inches(2.8), RGBColor(22,86,126))
    before_card.fill.fore_color.rgb = RGBColor(22,86,126)
    slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(4.2), Inches(0.4)).text_frame.text = "BEFORE"
    if os.path.exists(BEFORE_IMG):
        slide.shapes.add_picture(BEFORE_IMG, Inches(0.8), Inches(1.1), width=Inches(4.2), height=Inches(1.6))

    after_card = add_rounded_rect(slide, Inches(5.3), Inches(0.6), Inches(4.5), Inches(2.8), RGBColor(18,103,74))
    after_card.fill.fore_color.rgb = RGBColor(18,103,74)
    slide.shapes.add_textbox(Inches(5.5), Inches(0.7), Inches(4.2), Inches(0.4)).text_frame.text = "AFTER"
    if os.path.exists(AFTER_IMG):
        slide.shapes.add_picture(AFTER_IMG, Inches(5.5), Inches(1.1), width=Inches(4.2), height=Inches(1.6))

    # Action Plan bullets on bottom panel
    bottom = add_rounded_rect(slide, Inches(0.6), Inches(3.6), prs.slide_width - Inches(1.2), Inches(1.3), RGBColor(14,45,78))
    bottom.fill.fore_color.rgb = RGBColor(14,45,78)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(9.0), Inches(1.1))
    tf = tb.text_frame
    tf.text = "ACTION PLAN / RECOMMENDATIONS"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(210,230,250)
    p = tf.add_paragraph()
    p.text = "• Partner team to push pending code changes immediately\n• Reconnect ETL/schedulers and validate runs\n• Enable monitoring alerts and conduct joint validation\n• Sign-off migration after data flow stabilizes"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(230,240,255)


def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    panel = add_rounded_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RGBColor(4,22,40))
    panel.fill.fore_color.rgb = RGBColor(4,22,40)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(9.0), Inches(0.8))
    title.text_frame.text = "Conclusion & Next Steps"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 235, 250)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(9.0), Inches(2.0))
    tf = tb.text_frame
    tf.text = "Our team completed the migration work and provided full technical support. Current data staleness is caused by pending code/ETL reconnection on the consuming team's side. We've defined remediation steps and are ready to assist until final sign-off."
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(200,215,235)

    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(8.6), Inches(5.0), width=Inches(1.0))


def build_presentation():
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(10)
    add_title_slide(prs)
    add_executive_summary_slide(prs)
    add_migration_scope_slide(prs)
    add_support_provided_slide(prs)
    add_statistics_slide(prs, METRICS)
    add_issue_found_slide(prs)
    add_action_plan_slide(prs)
    add_closing_slide(prs)
    prs.save(OUT_FILE)
    print(f"Saved presentation to {OUT_FILE}")


if __name__ == "__main__":
    build_presentation()
